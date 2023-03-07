"""Slides parser.

Contains parsers for .pptx files.

"""

import os
from pathlib import Path
from typing import Dict

from gpt_index.readers.file.base_parser import BaseParser


class PptxParser(BaseParser):
    """Powerpoint parser.

    Extract text, caption images, and specify slides.

    """

    def _init_parser(self) -> Dict:
        """Init parser."""
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            raise ValueError(
                "The package `python-pptx` is required to read Powerpoint files."
            )
        try:
            import torch  # noqa: F401
        except ImportError:
            raise ValueError("The package `pytorch` is required to caption images.")
        try:
            from transformers import (
                AutoTokenizer,
                VisionEncoderDecoderModel,
                ViTFeatureExtractor,
            )
        except ImportError:
            raise ValueError(
                "The package `transformers` is required to caption images."
            )
        try:
            from PIL import Image  # noqa: F401
        except ImportError:
            raise ValueError(
                "PIL is required to read image files." "Please run `pip install Pillow`"
            )

        model = VisionEncoderDecoderModel.from_pretrained(
            "nlpconnect/vit-gpt2-image-captioning"
        )
        feature_extractor = ViTFeatureExtractor.from_pretrained(
            "nlpconnect/vit-gpt2-image-captioning"
        )
        tokenizer = AutoTokenizer.from_pretrained(
            "nlpconnect/vit-gpt2-image-captioning"
        )

        return {
            "feature_extractor": feature_extractor,
            "model": model,
            "tokenizer": tokenizer,
        }

    def caption_image(self, tmp_image_file: str) -> str:
        """Generate text caption of image."""
        import torch
        from PIL import Image

        model = self.parser_config["model"]
        feature_extractor = self.parser_config["feature_extractor"]
        tokenizer = self.parser_config["tokenizer"]

        device = "cuda" if torch.cuda.is_available() else "cpu"
        model.to(device)

        max_length = 16
        num_beams = 4
        gen_kwargs = {"max_length": max_length, "num_beams": num_beams}

        i_image = Image.open(tmp_image_file)
        if i_image.mode != "RGB":
            i_image = i_image.convert(mode="RGB")

        pixel_values = feature_extractor(
            images=[i_image], return_tensors="pt"
        ).pixel_values
        pixel_values = pixel_values.to(device)

        output_ids = model.generate(pixel_values, **gen_kwargs)

        preds = tokenizer.batch_decode(output_ids, skip_special_tokens=True)
        return preds[0].strip()

    def parse_file(self, file: Path, errors: str = "ignore") -> str:
        """Parse file."""
        from pptx import Presentation

        presentation = Presentation(file)
        result = ""
        for i, slide in enumerate(presentation.slides):
            result += f"\n\nSlide #{i}: \n"
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    # get image "file" contents
                    image_bytes = image.blob
                    # temporarily save the image to feed into model
                    image_filename = f"tmp_image.{image.ext}"
                    with open(image_filename, "wb") as f:
                        f.write(image_bytes)
                    result += f"\n Image: {self.caption_image(image_filename)}\n\n"

                    os.remove(image_filename)
                if hasattr(shape, "text"):
                    result += f"{shape.text}\n"

        return result
