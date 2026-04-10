"""
Default parsers for ConfluenceReader file types.
"""

import os
from io import BytesIO
from typing import Dict, Iterable, override

from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document

from .event import FileType


def _error_text(file_type: str, file_path: str, exc: Exception) -> str:
    file_name = os.path.basename(file_path)
    exception_message = f"{exc.__class__.__name__}: {exc}"
    return (
        f"error processing {file_type}; "
        f"exception message: {exception_message}; "
        f"file name: {file_name}"
    )


class DefaultPageHtmlParser(BaseReader):
    """Convert a Confluence page body (HTML) to Markdown using markdownify."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            from markdownify import markdownify  # type: ignore
        except ImportError:
            raise ImportError(
                "`markdownify` package not found. Install it directly: "
                "`pip install markdownify`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        with open(file_path, "r", encoding="utf-8") as f:
            html = f.read()
        if not html:
            return [Document(text="")]
        try:
            text = markdownify(
                html,
                heading_style="ATX",
                bullets="*",
                strip=["script", "style"],
            )
        except Exception as e:
            return [Document(text=_error_text("page_html", file_path, e))]
        return [Document(text=text.strip())]


class DefaultHtmlParser(BaseReader):
    """Extract plain text from an HTML file attachment using BeautifulSoup."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            from bs4 import BeautifulSoup  # type: ignore
        except ImportError:
            raise ImportError(
                "`beautifulsoup4` package not found. Install it directly: "
                "`pip install beautifulsoup4`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            soup = BeautifulSoup(content, "html.parser")
            text = soup.get_text(separator=" ", strip=True)
        except Exception as e:
            return [Document(text=_error_text("html", file_path, e))]
        return [Document(text=text.strip())]


class DefaultPdfParser(BaseReader):
    """Convert a PDF to text via pdf2image + pytesseract OCR."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import pytesseract  # type: ignore
            from pdf2image import convert_from_path  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract` or `pdf2image` package not found. Install them directly: "
                "`pip install pytesseract pdf2image`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            images = convert_from_path(file_path)
            text = ""
            for i, image in enumerate(images):
                text += f"Page {i + 1}:\n{pytesseract.image_to_string(image)}\n\n"
        except Exception as e:
            return [Document(text=_error_text("pdf", file_path, e))]
        return [Document(text=text.strip())]


class DefaultImageParser(BaseReader):
    """Extract text from an image file via pytesseract OCR."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract` or `Pillow` package not found. Install them directly: "
                "`pip install pytesseract Pillow`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            with Image.open(file_path) as image:
                text = pytesseract.image_to_string(image)
        except Exception as e:
            return [Document(text=_error_text("image", file_path, e))]
        return [Document(text=text.strip())]


class DefaultDocParser(BaseReader):
    """Extract text from a .docx file using docx2txt."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import docx2txt  # type: ignore
        except ImportError:
            raise ImportError(
                "`docx2txt` package not found. Install it directly: "
                "`pip install docx2txt`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            text = docx2txt.process(file_path)
        except Exception as e:
            return [Document(text=_error_text("document", file_path, e))]
        return [Document(text=text or "")]


class DefaultPptParser(BaseReader):
    """Extract text from a PowerPoint file using python-pptx."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            raise ImportError(
                "`python-pptx` package not found. Install it directly: "
                "`pip install python-pptx`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
        except Exception as e:
            return [Document(text=_error_text("presentation", file_path, e))]
        return [Document(text=text.strip())]


class DefaultXlsParser(BaseReader):
    """Extract text from an Excel (.xlsx/.xls) file using pandas + openpyxl."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import pandas as pd  # type: ignore
        except ImportError:
            raise ImportError(
                "`pandas` package not found. Install it directly: "
                "`pip install pandas openpyxl`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
        except Exception as e:
            return [Document(text=_error_text("spreadsheet", file_path, e))]
        text = ""
        for sheet_name, sheet_data in sheets.items():
            text += f"{sheet_name}:\n"
            for _, row in sheet_data.iterrows():
                text += "\t".join(str(value) for value in row) + "\n"
            text += "\n"
        return [Document(text=text.strip())]


class DefaultXlsbParser(BaseReader):
    """Extract text from an Excel Binary (.xlsb) file using pandas + pyxlsb."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import pandas as pd  # type: ignore
        except ImportError:
            raise ImportError(
                "`pandas` package not found. Install it directly: "
                "`pip install pandas pyxlsb`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            df = pd.read_excel(file_path, engine="pyxlsb")
            text_rows = [", ".join(row.astype(str)) for _, row in df.iterrows()]
            text = "\n".join(text_rows)
        except Exception as e:
            return [Document(text=_error_text("xlsb", file_path, e))]
        return [Document(text=text.strip())]


class DefaultCsvParser(BaseReader):
    """Extract text from a CSV file using pandas."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import pandas as pd  # type: ignore
        except ImportError:
            raise ImportError(
                "`pandas` package not found. Install it directly: "
                "`pip install pandas`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            df = pd.read_csv(file_path, low_memory=False)
            text_rows = [", ".join(row.astype(str)) for _, row in df.iterrows()]
            text = "\n".join(text_rows)
        except Exception as e:
            return [Document(text=_error_text("csv", file_path, e))]
        return [Document(text=text.strip())]


class DefaultTxtParser(BaseReader):
    """Read plain text files.  No 3rd-party dependencies."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
        except Exception as e:
            return [Document(text=_error_text("text", file_path, e))]
        return [Document(text=text.strip())]


class DefaultMsgParser(BaseReader):
    """Extract text from an Outlook .msg file using extract-msg."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import extract_msg  # type: ignore
        except ImportError:
            raise ImportError(
                "`extract-msg` package not found. Install it directly: "
                "`pip install extract-msg`, or install all parsers: "
                '`pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            with extract_msg.Message(file_path) as msg:
                text = (
                    f"Subject: {msg.subject}\nFrom: {msg.sender}\n"
                    f"To: {msg.to}\nCC: {msg.cc}\n\n{msg.body}"
                )
        except Exception as e:
            return [Document(text=_error_text("msg", file_path, e))]
        return [Document(text=text.strip())]


class DefaultSvgParser(BaseReader):
    """Render an SVG to a raster image then OCR it via pytesseract."""

    @override
    def lazy_load_data(self, file_path: str, **kwargs) -> Iterable[Document]:
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            from reportlab.graphics import renderPM  # type: ignore
            from svglib.svglib import svg2rlg  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract`, `Pillow`, `svglib`, or `reportlab` package not found. "
                "Install them directly: `pip install pytesseract Pillow svglib reportlab`, "
                'or install all parsers: `pip install "llama-index-readers-confluence[all]"`'
            )
        try:
            drawing = svg2rlg(file_path)
            if drawing is None:
                return [Document(text="")]
            img_data = BytesIO()
            renderPM.drawToFile(drawing, img_data, fmt="PNG")
            img_data.seek(0)
            with Image.open(img_data) as image:
                text = pytesseract.image_to_string(image)
        except Exception as e:
            return [Document(text=_error_text("svg", file_path, e))]
        return [Document(text=text.strip())]


def get_default_parsers() -> Dict[FileType, BaseReader]:
    """
    Return the default parser mapping used by ConfluenceReader.
    """
    return {
        FileType.PAGE_HTML: DefaultPageHtmlParser(),
        FileType.HTML: DefaultHtmlParser(),
        FileType.PDF: DefaultPdfParser(),
        FileType.IMAGE: DefaultImageParser(),
        FileType.DOCUMENT: DefaultDocParser(),
        FileType.PRESENTATION: DefaultPptParser(),
        FileType.SPREADSHEET: DefaultXlsParser(),
        FileType.XLSB: DefaultXlsbParser(),
        FileType.CSV: DefaultCsvParser(),
        FileType.TEXT: DefaultTxtParser(),
        FileType.MSG: DefaultMsgParser(),
        FileType.SVG: DefaultSvgParser(),
    }
