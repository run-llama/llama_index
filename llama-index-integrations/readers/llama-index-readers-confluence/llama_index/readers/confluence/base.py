"""Confluence reader."""

import logging
import os
import uuid
import tempfile
from typing import Callable, Dict, List, Optional
from urllib.parse import unquote

import requests
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.core.instrumentation import DispatcherSpanMixin, get_dispatcher
from retrying import retry
from io import BytesIO

from .event import (
    FileType,
    TotalPagesToProcessEvent,
    PageDataFetchStartedEvent,
    PageDataFetchCompletedEvent,
    PageSkippedEvent,
    PageFailedEvent,
    AttachmentProcessingStartedEvent,
    AttachmentProcessedEvent,
    AttachmentSkippedEvent,
    AttachmentFailedEvent,
)

CONFLUENCE_API_TOKEN = "CONFLUENCE_API_TOKEN"
CONFLUENCE_PASSWORD = "CONFLUENCE_PASSWORD"
CONFLUENCE_USERNAME = "CONFLUENCE_USERNAME"

internal_logger = logging.getLogger(__name__)
dispatcher = get_dispatcher(__name__)


class CustomParserManager:
    def __init__(
        self, custom_parsers: Optional[Dict[FileType, BaseReader]], custom_folder: str
    ):
        self.custom_parsers = custom_parsers or {}
        self.custom_folder = custom_folder

    def __remove_custom_file(self, file_path: str):
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            print(f"Error removing file {file_path}: {e}")

    def process_with_custom_parser(
        self, file_type: FileType, file_content: bytes, extension: str
    ) -> Optional[str]:
        if file_type not in self.custom_parsers:
            return None

        file_name = f"{uuid.uuid4().hex}.{extension}"
        custom_file_path = os.path.join(self.custom_folder, file_name)
        with open(custom_file_path, "wb") as f:
            f.write(file_content)

        try:
            markdown_text = "\n".join(
                doc.text
                for doc in self.custom_parsers[file_type].load_data(
                    file_path=custom_file_path
                )
            )
        finally:
            self.__remove_custom_file(custom_file_path)
        return markdown_text


class ConfluenceReader(BaseReader, DispatcherSpanMixin):
    """
    Confluence reader.

    Reads a set of confluence pages given a space key and optionally a list of page ids

    For more on OAuth login, checkout:
        - https://atlassian-python-api.readthedocs.io/index.html
        - https://developer.atlassian.com/cloud/confluence/oauth-2-3lo-apps/

    Args:
        oauth2 (dict): Atlassian OAuth 2.0, minimum fields are `client_id` and `token`, where `token` is a dict and must at least contain "access_token" and "token_type".
        base_url (str): 'base_url' for confluence cloud instance, this is suffixed with '/wiki', eg 'https://yoursite.atlassian.com/wiki'
        cloud (bool): connecting to Confluence Cloud or self-hosted instance
        api_token (str): Confluence API token, see https://confluence.atlassian.com/cloud/api-tokens-938839638.html
        cookies (dict): Confluence cookies, see https://atlassian-python-api.readthedocs.io/index.html
        user_name (str): Confluence username, used for basic auth. Must be used with `password`.
        password (str): Confluence password, used for basic auth. Must be used with `user_name`.
        client_args (dict): Additional keyword arguments to pass directly to the Atlassian Confluence client constructor, for example `{'backoff_and_retry': True}`.
        custom_parsers (dict): Optional custom parsers for different file types. Maps FileType enum values to BaseReader instances.
        process_attachment_callback (callable): Optional callback function to determine whether to process an attachment. Should return tuple[bool, str] where bool indicates whether to process and str provides reason if not.
        process_document_callback (callable): Optional callback function to determine whether to process a document. Should return bool indicating whether to process.
        custom_folder (str): Optional custom folder path for storing temporary files. Can only be used when custom_parsers are provided. Defaults to current working directory if custom_parsers are used.
        logger (Logger): Optional custom logger instance. If not provided, uses internal logger.
        fail_on_error (bool): Whether to raise exceptions on processing errors or continue with warnings. Default is True.

    Instrumentation Events:
        The ConfluenceReader uses LlamaIndex's instrumentation system to emit events during document and attachment processing.
        These events can be captured by adding event handlers to the dispatcher.

        Available events:
        - TotalPagesToProcessEvent: Emitted when the total number of pages to process is determined
        - PageDataFetchStartedEvent: Emitted when processing of a page begins
        - PageDataFetchCompletedEvent: Emitted when a page is successfully processed
        - PageFailedEvent: Emitted when page processing fails
        - PageSkippedEvent: Emitted when a page is skipped due to callback decision
        - AttachmentProcessingStartedEvent: Emitted when attachment processing begins
        - AttachmentProcessedEvent: Emitted when an attachment is successfully processed
        - AttachmentSkippedEvent: Emitted when an attachment is skipped
        - AttachmentFailedEvent: Emitted when attachment processing fails

        To listen to events, add an event handler to the dispatcher:
        ```python
        from llama_index.core.instrumentation import get_dispatcher
        from llama_index.core.instrumentation.event_handlers import BaseEventHandler

        class MyEventHandler(BaseEventHandler):
            def handle(self, event):
                print(f"Event: {event.class_name()}")

        dispatcher = get_dispatcher(__name__)
        dispatcher.add_event_handler(MyEventHandler())
        ```

    """

    def __init__(
        self,
        base_url: str = None,
        oauth2: Optional[Dict] = None,
        cloud: bool = True,
        api_token: Optional[str] = None,
        cookies: Optional[dict] = None,
        user_name: Optional[str] = None,
        password: Optional[str] = None,
        client_args: Optional[dict] = None,
        custom_parsers: Optional[Dict[FileType, BaseReader]] = None,
        process_attachment_callback: Optional[
            Callable[[str, int], tuple[bool, str]]
        ] = None,
        process_document_callback: Optional[Callable[[str], bool]] = None,
        custom_folder: Optional[str] = None,
        logger: Optional[logging.Logger] = None,
        fail_on_error: bool = True,
    ) -> None:
        if base_url is None:
            raise ValueError("Must provide `base_url`")

        self.base_url = base_url
        self.custom_parsers = custom_parsers or {}

        # Only set custom_folder if custom_parsers are provided
        if custom_parsers and custom_folder:
            self.custom_folder = custom_folder
        elif custom_parsers:
            self.custom_folder = os.getcwd()
        elif custom_folder:
            raise ValueError(
                "custom_folder can only be used when custom_parsers are provided"
            )
        else:
            self.custom_folder = None

        self.logger = logger or internal_logger
        self.process_attachment_callback = process_attachment_callback
        self.process_document_callback = process_document_callback
        self.fail_on_error = fail_on_error

        try:
            from atlassian import Confluence
        except ImportError:
            raise ImportError(
                "`atlassian` package not found, please run `pip install atlassian-python-api`"
            )
        self.confluence: Confluence = None
        if client_args is None:
            client_args = {}
        if oauth2:
            self.confluence = Confluence(
                url=base_url, oauth2=oauth2, cloud=cloud, **client_args
            )
        else:
            if api_token is not None:
                self.confluence = Confluence(
                    url=base_url, token=api_token, cloud=cloud, **client_args
                )
            elif cookies is not None:
                self.confluence = Confluence(
                    url=base_url, cookies=cookies, cloud=cloud, **client_args
                )
            elif user_name is not None and password is not None:
                self.confluence = Confluence(
                    url=base_url,
                    username=user_name,
                    password=password,
                    cloud=cloud,
                    **client_args,
                )
            else:
                api_token = os.getenv(CONFLUENCE_API_TOKEN)
                if api_token is not None:
                    self.confluence = Confluence(
                        url=base_url, token=api_token, cloud=cloud, **client_args
                    )
                else:
                    user_name = os.getenv(CONFLUENCE_USERNAME)
                    password = os.getenv(CONFLUENCE_PASSWORD)
                    if user_name is not None and password is not None:
                        self.confluence = Confluence(
                            url=base_url,
                            username=user_name,
                            password=password,
                            cloud=cloud,
                            **client_args,
                        )
                    else:
                        raise ValueError(
                            "Must set one of environment variables `CONFLUENCE_API_KEY`, or"
                            " `CONFLUENCE_USERNAME` and `CONFLUENCE_PASSWORD`, if oauth2, or"
                            " api_token, or user_name and password parameters are not provided"
                        )

        self._next_cursor = None
        if custom_parsers:
            self.custom_parser_manager = CustomParserManager(
                custom_parsers, self.custom_folder
            )
        else:
            self.custom_parser_manager = None

    def _format_attachment_header(self, attachment: dict) -> str:
        """Formats the attachment title as a markdown header."""
        return f"# {attachment['title']}\n"

    @dispatcher.span
    def load_data(
        self,
        space_key: Optional[str] = None,
        page_ids: Optional[List[str]] = None,
        folder_id: Optional[str] = None,
        page_status: Optional[str] = None,
        label: Optional[str] = None,
        cql: Optional[str] = None,
        include_attachments=False,
        include_children=False,
        start: Optional[int] = None,
        cursor: Optional[str] = None,
        limit: Optional[int] = None,
        max_num_results: Optional[int] = None,
    ) -> List[Document]:
        """
        Load Confluence pages from Confluence, specifying by one of four mutually exclusive methods:
        `space_key`, `page_ids`, `label`, `folder_id` or `cql`
        (Confluence Query Language https://developer.atlassian.com/cloud/confluence/advanced-searching-using-cql/ ).

        Args:
            space_key (str): Confluence space key, eg 'DS'
            page_ids (list): List of page ids, eg ['123456', '123457']
            folder_id (str): Confluence folder id, eg '1234567890'
            page_status (str): Page status, one of None (all statuses), 'current', 'draft', 'archived'.  Only compatible with space_key.
            label (str): Confluence label, eg 'my-label'
            cql (str): Confluence Query Language query, eg 'label="my-label"'
            include_attachments (bool): If True, include attachments.
            include_children (bool): If True, do a DFS of the descendants of each page_id in `page_ids`.  Only compatible with `page_ids`.
            start (int): Skips over the first n elements. Used only with space_key
            cursor (str): Skips to the cursor. Used with cql and label, set when the max limit has been hit for cql based search
            limit (int): Deprecated, use `max_num_results` instead.
            max_num_results (int): Maximum number of results to return.  If None, return all results.  Requests are made in batches to achieve the desired number of results.

        """
        num_space_key_parameter = 1 if space_key else 0
        num_page_ids_parameter = 1 if page_ids is not None else 0
        num_label_parameter = 1 if label else 0
        num_cql_parameter = 1 if cql else 0
        num_folder_id_parameter = 1 if folder_id else 0
        if (
            num_space_key_parameter
            + num_page_ids_parameter
            + num_label_parameter
            + num_cql_parameter
            + num_folder_id_parameter
            != 1
        ):
            raise ValueError(
                "Must specify exactly one among `space_key`, `page_ids`, `label`, `cql`"
                " parameters."
            )

        if cursor and start:
            raise ValueError("Must not specify `start` when `cursor` is specified")

        if space_key and cursor:
            raise ValueError("Must not specify `cursor` when `space_key` is specified")

        if page_status and not space_key:
            raise ValueError(
                "Must specify `space_key` when `page_status` is specified."
            )

        if include_children and not page_ids:
            raise ValueError(
                "Must specify `page_ids` when `include_children` is specified."
            )

        if limit is not None:
            max_num_results = limit
            self.logger.warning(
                "`limit` is deprecated and no longer relates to the Confluence server's"
                " API limits.  If you wish to limit the number of returned results"
                " please use `max_num_results` instead."
            )

        try:
            import html2text  # type: ignore
        except ImportError:
            raise ImportError(
                "`html2text` package not found, please run `pip install html2text`"
            )

        text_maker = html2text.HTML2Text()

        if not start:
            start = 0

        pages: List = []
        if space_key:
            pages.extend(
                self._get_data_with_paging(
                    self.confluence.get_all_pages_from_space,
                    start=start,
                    max_num_results=max_num_results,
                    space=space_key,
                    status=page_status,
                    expand="body.export_view.value",
                    content_type="page",
                )
            )
        elif label:
            pages.extend(
                self._get_cql_data_with_paging(
                    start=start,
                    cursor=cursor,
                    cql=f'type="page" AND label="{label}"',
                    max_num_results=max_num_results,
                    expand="body.export_view.value",
                )
            )
        elif cql:
            pages.extend(
                self._get_cql_data_with_paging(
                    start=start,
                    cursor=cursor,
                    cql=cql,
                    max_num_results=max_num_results,
                    expand="body.export_view.value",
                )
            )
        elif page_ids:
            if include_children:
                dfs_page_ids = []
                max_num_remaining = max_num_results
                for page_id in page_ids:
                    current_dfs_page_ids = self._dfs_page_ids(
                        page_id,
                        type="page",
                        max_num_results=max_num_remaining,
                    )
                    dfs_page_ids.extend(current_dfs_page_ids)
                    if max_num_results is not None:
                        max_num_remaining -= len(current_dfs_page_ids)
                        if max_num_remaining <= 0:
                            break
                page_ids = dfs_page_ids
            for page_id in (
                page_ids[:max_num_results] if max_num_results is not None else page_ids
            ):
                pages.append(
                    self._get_data_with_retry(
                        self.confluence.get_page_by_id,
                        page_id=page_id,
                        expand="body.export_view.value",
                    )
                )
        elif folder_id:
            # Fetch all folders in the folder
            max_num_remaining = max_num_results
            page_ids = self._dfs_page_ids(
                folder_id,
                type="folder",
                max_num_results=max_num_remaining,
            )
            for page_id in (
                page_ids[:max_num_results] if max_num_results is not None else page_ids
            ):
                pages.append(
                    self._get_data_with_retry(
                        self.confluence.get_page_by_id,
                        page_id=page_id,
                        expand="body.export_view.value",
                    )
                )

        docs = []

        if pages:
            dispatcher.event(TotalPagesToProcessEvent(total_pages=len(pages)))

        for page in pages:
            try:
                doc = self.process_page(page, include_attachments, text_maker)
                if doc:
                    docs.append(doc)
            except Exception as e:
                self.logger.error(f"Error processing page: {e}")
                dispatcher.event(PageFailedEvent(page_id=page["id"], error=str(e)))
                if self.fail_on_error:
                    raise
                else:
                    self.logger.warning(
                        f"Failed to process page {page['id']}: {e}. Skipping this page."
                    )
                    continue
        return docs

    def _dfs_page_ids(self, id, type="page", max_num_results=None):
        ret = []

        max_num_remaining = max_num_results
        if type == "page":
            ret.append(id)
            if max_num_remaining is not None:
                max_num_remaining -= 1
                if max_num_remaining < 0:
                    return ret

        # Fetch both page and folder children with their types
        child_items = [
            (child_id, "page")
            for child_id in self._get_data_with_paging(
                self.confluence.get_child_id_list,
                page_id=id,
                type="page",
                max_num_results=max_num_remaining,
            )
        ] + [
            (child_id, "folder")
            for child_id in self._get_data_with_paging(
                self.confluence.get_child_id_list,
                page_id=id,
                type="folder",
                max_num_results=max_num_remaining,
            )
        ]

        for child_id, child_type in child_items:
            if max_num_remaining is not None and max_num_remaining <= 0:
                break

            dfs_ids = self._dfs_page_ids(
                child_id, type=child_type, max_num_results=max_num_remaining
            )
            ret.extend(dfs_ids)

            if max_num_remaining is not None:
                max_num_remaining -= len(dfs_ids)
                if max_num_remaining <= 0:
                    break

        return ret

    def _get_data_with_paging(
        self, paged_function, start=0, max_num_results=50, **kwargs
    ):
        max_num_remaining = max_num_results
        ret = []
        while True:
            results = self._get_data_with_retry(
                paged_function, start=start, limit=max_num_remaining, **kwargs
            )
            ret.extend(results)
            if (
                len(results) == 0
                or max_num_results is not None
                and len(results) >= max_num_remaining
            ):
                break

            start += len(results)
            if max_num_remaining is not None:
                max_num_remaining -= len(results)
        return ret

    def _get_cql_data_with_paging(
        self,
        cql,
        start=0,
        cursor=None,
        max_num_results=50,
        expand="body.export_view.value",
    ):
        max_num_remaining = max_num_results
        ret = []
        params = {"cql": cql, "start": start, "expand": expand}
        if cursor:
            params["cursor"] = unquote(cursor)

        if max_num_results is not None:
            params["limit"] = max_num_remaining
        while True:
            results = self._get_data_with_retry(
                self.confluence.get, path="rest/api/content/search", params=params
            )
            ret.extend(results["results"])

            params["start"] += len(results["results"])

            next_url = (
                results["_links"]["next"] if "next" in results["_links"] else None
            )
            if not next_url:
                self._next_cursor = None
                break

            if "cursor=" in next_url:  # On confluence Server this is not set
                cursor = next_url.split("cursor=")[1].split("&")[0]
                params["cursor"] = unquote(cursor)

            if max_num_results is not None:
                params["limit"] -= len(results["results"])
                if params["limit"] <= 0:
                    self._next_cursor = cursor
                    break

        return ret

    def get_next_cursor(self):
        """
        Returns: The last set cursor from a cql based search.
        """
        return self._next_cursor

    @retry(stop_max_attempt_number=1, wait_fixed=4)
    def _get_data_with_retry(self, function, **kwargs):
        return function(**kwargs)

    @dispatcher.span
    def process_page(self, page, include_attachments, text_maker):
        self.logger.info("Processing " + self.base_url + page["_links"]["webui"])

        if self.process_document_callback:
            should_process = self.process_document_callback(page["id"])
            if not should_process:
                self.logger.info(
                    f"Skipping page {page['id']} based on callback decision."
                )
                dispatcher.event(PageSkippedEvent(page_id=page["id"]))
                return None

        dispatcher.event(PageDataFetchStartedEvent(page_id=page["id"]))

        if include_attachments:
            attachment_texts = self.process_attachment(page["id"])
        else:
            attachment_texts = []
        if FileType.HTML in self.custom_parsers and self.custom_folder:
            html_text = page["body"]["export_view"]["value"]
            # save in temporary file
            file_location = None
            with tempfile.NamedTemporaryFile(
                mode="w",
                suffix=".html",
                encoding="utf-8",
                dir=self.custom_folder,
                delete=False,
            ) as f:
                f.write(html_text)
                file_location = f.name
            try:
                text = (
                    page["title"]
                    + "\n"
                    + "\n".join(
                        doc.text
                        for doc in self.custom_parsers[FileType.HTML].load_data(
                            file_path=file_location
                        )
                    )
                    + "\n"
                    + "\n".join(attachment_texts)
                )
            finally:
                try:
                    os.unlink(file_location)
                except OSError:
                    pass
        else:
            text = text_maker.handle(page["body"]["export_view"]["value"]) + "".join(
                attachment_texts
            )

        doc = Document(
            text=text,
            doc_id=page["id"],
            extra_info={
                "title": page["title"],
                "page_id": page["id"],
                "status": page["status"],
                "url": self.base_url + page["_links"]["webui"],
            },
        )
        dispatcher.event(PageDataFetchCompletedEvent(page_id=page["id"], document=doc))
        return doc

    @dispatcher.span
    def process_attachment(self, page_id):
        try:
            pass
        except ImportError:
            raise ImportError(
                "`pytesseract` or `pdf2image` or `Pillow` package not found, please run"
                " `pip install pytesseract pdf2image Pillow`"
            )

        # depending on setup you may also need to set the correct path for poppler and tesseract
        attachments = self.confluence.get_attachments_from_content(page_id)["results"]
        texts = []
        if not attachments:
            return texts

        for attachment in attachments:
            self.logger.info("Processing attachment " + attachment["title"])
            dispatcher.event(
                AttachmentProcessingStartedEvent(
                    page_id=page_id,
                    attachment_id=attachment["id"],
                    attachment_name=attachment["title"],
                    attachment_type=attachment["metadata"]["mediaType"],
                    attachment_size=attachment["extensions"]["fileSize"],
                    attachment_link=attachment["_links"]["webui"],
                )
            )

            if self.process_attachment_callback:
                should_process, reason = self.process_attachment_callback(
                    attachment["metadata"]["mediaType"],
                    attachment["extensions"]["fileSize"],
                    attachment["title"],
                )
                if not should_process:
                    self.logger.info(
                        f"Skipping attachment {attachment['title']} based on callback decision."
                    )
                    dispatcher.event(
                        AttachmentSkippedEvent(
                            page_id=page_id,
                            attachment_id=attachment["id"],
                            attachment_name=attachment["title"],
                            attachment_type=attachment["metadata"]["mediaType"],
                            attachment_size=attachment["extensions"]["fileSize"],
                            attachment_link=attachment["_links"]["webui"],
                            reason=reason,
                        )
                    )
                    continue

            try:
                media_type = attachment["metadata"]["mediaType"]
                absolute_url = self.base_url + attachment["_links"]["download"]
                title = self._format_attachment_header(attachment)
                if media_type == "application/pdf":
                    self.logger.info("Processing PDF attachment " + absolute_url)
                    text = title + self.process_pdf(absolute_url)
                elif (
                    media_type == "image/png"
                    or media_type == "image/jpg"
                    or media_type == "image/jpeg"
                    or media_type == "image/webp"
                ):
                    self.logger.info("Processing image attachment " + absolute_url)
                    text = title + self.process_image(absolute_url)
                elif (
                    media_type
                    == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ):
                    self.logger.info(
                        "Processing Word document attachment " + absolute_url
                    )
                    text = title + self.process_doc(absolute_url)
                elif (
                    media_type == "application/vnd.ms-excel"
                    or media_type
                    == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    or media_type == "application/vnd.ms-excel.sheet.macroenabled.12"
                ):
                    if attachment["title"].endswith(".csv") or absolute_url.endswith(
                        ".csv"
                    ):
                        self.logger.info("Processing CSV attachment " + absolute_url)
                        text = title + self.process_csv(absolute_url)
                    else:
                        self.logger.info("Processing XLS attachment " + absolute_url)
                        text = title + self.process_xls(absolute_url)
                elif (
                    media_type
                    == "application/vnd.ms-excel.sheet.binary.macroenabled.12"
                ):
                    self.logger.info("Processing XLSB attachment " + absolute_url)
                    text = title + self.process_xlsb(absolute_url)
                elif media_type == "text/csv":
                    self.logger.info("Processing CSV attachment " + absolute_url)
                    text = title + self.process_csv(absolute_url)
                elif media_type == "application/vnd.ms-outlook":
                    self.logger.info(
                        "Processing Outlook message attachment " + absolute_url
                    )
                    text = title + self.process_msg(absolute_url)
                elif media_type == "text/html":
                    self.logger.info("  Processing HTML attachment " + absolute_url)
                    text = title + self.process_html(absolute_url)
                elif media_type == "text/plain":
                    if attachment["title"].endswith(".csv") or absolute_url.endswith(
                        ".csv"
                    ):
                        self.logger.info("Processing CSV attachment " + absolute_url)
                        text = title + self.process_csv(absolute_url)
                    else:
                        self.logger.info("Processing Text attachment " + absolute_url)
                        text = title + self.process_txt(absolute_url)
                elif media_type == "text/markdown" or absolute_url.endswith(
                    (".md", ".mdx")
                ):
                    self.logger.info("Processing Markdown attachment " + absolute_url)
                    text = title + self.process_txt(absolute_url)
                elif media_type == "image/svg+xml":
                    self.logger.info("Processing SVG attachment " + absolute_url)
                    text = title + self.process_svg(absolute_url)
                elif (
                    media_type
                    == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    or media_type
                    == "application/vnd.ms-powerpoint.presentation.macroenabled.12"
                ):
                    self.logger.info(
                        "Processing PowerPoint attachment "
                        + absolute_url
                        + " ("
                        + media_type
                        + ")"
                    )
                    text = title + self.process_ppt(absolute_url)
                elif media_type == "binary/octet-stream" and (
                    attachment["title"].strip().endswith(".md")
                    or attachment["title"].strip().endswith(".mdx")
                ):
                    self.logger.info("Processing Markdown attachment " + absolute_url)
                    text = title + self.process_txt(absolute_url)
                else:
                    self.logger.info(
                        f"Skipping unsupported attachment {absolute_url} of media_type {media_type}"
                    )
                    dispatcher.event(
                        AttachmentSkippedEvent(
                            page_id=page_id,
                            attachment_id=attachment["id"],
                            attachment_name=attachment["title"],
                            attachment_type=attachment["metadata"]["mediaType"],
                            attachment_size=attachment["extensions"]["fileSize"],
                            attachment_link=attachment["_links"]["webui"],
                            reason="Unsupported media type",
                        )
                    )
                    continue
                texts.append(text)
                dispatcher.event(
                    AttachmentProcessedEvent(
                        page_id=page_id,
                        attachment_id=attachment["id"],
                        attachment_name=attachment["title"],
                        attachment_type=attachment["metadata"]["mediaType"],
                        attachment_size=attachment["extensions"]["fileSize"],
                        attachment_link=attachment["_links"]["webui"],
                    )
                )
            except Exception as e:
                self.logger.error(
                    f"Failed to process attachment {attachment['title']}: {e}"
                )
                dispatcher.event(
                    AttachmentFailedEvent(
                        page_id=page_id,
                        attachment_id=attachment["id"],
                        attachment_name=attachment["title"],
                        attachment_type=attachment["metadata"]["mediaType"],
                        attachment_size=attachment["extensions"]["fileSize"],
                        attachment_link=attachment["_links"]["webui"],
                        error=str(e),
                    )
                )

        return texts

    def process_pdf(self, link):
        if FileType.PDF not in self.custom_parsers:
            try:
                import pytesseract  # type: ignore
                from pdf2image import convert_from_bytes  # type: ignore
            except ImportError:
                raise ImportError(
                    "`pytesseract` or `pdf2image` package not found, please run `pip"
                    " install pytesseract pdf2image`"
                )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        if FileType.PDF in self.custom_parsers and self.custom_parser_manager:
            return self.custom_parser_manager.process_with_custom_parser(
                FileType.PDF, response.content, "pdf"
            )

        try:
            images = convert_from_bytes(response.content)
        except ValueError:
            return text

        for i, image in enumerate(images):
            image_text = pytesseract.image_to_string(image)
            text += f"Page {i + 1}:\n{image_text}\n\n"

        return text

    def process_html(self, link):
        try:
            from bs4 import BeautifulSoup  # type: ignore
        except ImportError:
            raise ImportError(
                "`beautifulsoup4` or `requests` package not found, please run `pip install beautifulsoup4 requests`"
            )

        try:
            response = self.confluence.request(path=link, absolute=True)
            if response.status_code != 200:
                return "Error fetching HTML content: HTTP Status Code {}".format(
                    response.status_code
                )

            if FileType.HTML in self.custom_parsers and self.custom_parser_manager:
                return self.custom_parser_manager.process_with_custom_parser(
                    FileType.HTML, response.content, "html"
                )

            # Parse the HTML content and extract text
            soup = BeautifulSoup(response.content, "html.parser")
            return soup.get_text(separator=" ", strip=True)
        except Exception as e:
            self.logger.error(f"Error processing HTML file at {link}: {e}")
            return f"Error processing HTML file: {link}. An error occurred while fetching or parsing the content."

    def process_txt(self, link):
        try:
            response = self.confluence.request(path=link, absolute=True)
            if response.status_code != 200:
                return "Error fetching text content: HTTP Status Code {}".format(
                    response.status_code
                )
            return response.text
        except Exception as e:
            self.logger.error(f"Error processing text file at {link}: {e}")
            return f"Error processing text file: {link}. An error occurred while fetching the content."

    def process_msg(self, link):
        try:
            import extract_msg  # type: ignore
        except ImportError:
            raise ImportError(
                "`extract-msg` package not found, please run `pip install extract-msg`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if response.status_code != 200 or response.content in [b"", None]:
            self.logger.error(f"Failed to download .msg file from {link}")
            return text

        file_data = BytesIO(response.content)

        try:
            # Load the .msg file content
            with extract_msg.Message(file_data) as msg:
                subject = msg.subject
                sender = msg.sender
                to = msg.to
                cc = msg.cc
                body = msg.body

                # Compile the extracted information into a text string
                text = (
                    f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nCC: {cc}\n\n{body}"
                )
        except Exception as e:
            self.logger.error(f"Error processing .msg file at {link}: {e}")
            return "Error processing .msg file."

        return text

    def process_image(self, link):
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract` or `Pillow` package not found, please run `pip install"
                " pytesseract Pillow`"
            )

        text = ""

        try:
            response = self.confluence.request(path=link, absolute=True)
            # Check if the response status code indicates success (200 OK)
            if response.status_code == 200 and response.content:
                try:
                    image = Image.open(BytesIO(response.content))
                    text = pytesseract.image_to_string(image)
                except OSError:
                    # Handle errors that occur while opening or processing the image
                    self.logger.error(
                        f"Error processing image at {link}: Unable to open or read the image content."
                    )
                    return text
            else:
                # Log non-200 responses here if needed
                self.logger.error(
                    f"Error fetching image at {link}: HTTP status code {response.status_code}."
                )
                return text
        except requests.exceptions.RequestException as e:
            # This catches any Requests-related exceptions, including HTTPError, ConnectionError, etc.
            self.logger.error(f"Request error while fetching image at {link}: {e}")
            return text

        return text

    def process_doc(self, link):
        try:
            import zipfile  # Import zipfile to catch BadZipFile exceptions
        except ImportError:
            raise ImportError("Failed to import BytesIO from io")
        if not self.custom_parsers.get(FileType.DOCUMENT):
            try:
                import docx2txt
            except ImportError:
                raise ImportError(
                    "`docx2txt` package not found, please run `pip install docx2txt`"
                )

        text = ""

        try:
            response = self.confluence.request(path=link, absolute=True)
            if response.status_code != 200 or response.content in [b"", None]:
                self.logger.error(
                    f"Error fetching document at {link}: HTTP status code {response.status_code}."
                )
                return text

            file_data = BytesIO(response.content)

            # save in file
            # Use custom parser if available
            if FileType.DOCUMENT in self.custom_parsers and self.custom_parser_manager:
                return self.custom_parser_manager.process_with_custom_parser(
                    FileType.DOCUMENT, file_data.getbuffer(), "docx"
                )

            try:
                text = docx2txt.process(file_data)
            except zipfile.BadZipFile:
                self.logger.error(
                    f"Error processing Word document at {link}: File is not a zip file."
                )
                return text
        except Exception as e:
            self.logger.error(f"Unexpected error processing document at {link}: {e}")
            return text

        return text

    def process_ppt(self, link):
        if not self.custom_parsers.get(FileType.PRESENTATION):
            try:
                from pptx import Presentation  # type: ignore
            except ImportError:
                raise ImportError(
                    "`python-pptx` package not found, please run `pip install python-pptx`"
                )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        if not file_data:
            self.logger.error(
                f"Error processing PowerPoint file at {link}: Empty content."
            )
            return text

        if FileType.PRESENTATION in self.custom_parsers and self.custom_parser_manager:
            return self.custom_parser_manager.process_with_custom_parser(
                FileType.PRESENTATION, file_data.getbuffer(), "pptx"
            )

        # Check if the response content is empty

        try:
            presentation = Presentation(file_data)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
        except (
            Exception
        ) as e:  # Catching a general exception to handle any unexpected errors
            self.logger.error(f"Error processing PowerPoint file at {link}: {e}")
            text = f"Error processing PowerPoint file: {link}. The file might be corrupt or not a valid PowerPoint file."

        return text.strip()  # Remove any leading/trailing whitespace

    def process_xls(self, link):
        try:
            import pandas as pd  # type: ignore
        except ImportError:
            raise ImportError(
                "`pandas` package not found, please run `pip install pandas`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        if FileType.SPREADSHEET in self.custom_parsers and self.custom_parser_manager:
            return self.custom_parser_manager.process_with_custom_parser(
                FileType.SPREADSHEET, file_data.getbuffer(), "xlsx"
            )
        # Try to read the Excel file
        try:
            # Use pandas to read all sheets; returns a dict of DataFrame
            sheets = pd.read_excel(file_data, sheet_name=None, engine="openpyxl")
        except Exception as e:
            return f"Failed to read Excel file: {e!s}"

        for sheet_name, sheet_data in sheets.items():
            text += f"{sheet_name}:\n"
            for row_index, row in sheet_data.iterrows():
                text += "\t".join(str(value) for value in row) + "\n"
            text += "\n"

        return text.strip()

    def process_xlsb(self, link):
        try:
            import pandas as pd
        except ImportError:
            raise ImportError(
                "`pandas` package not found, please run `pip install pandas`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        try:
            # Use pandas to read the .xlsb file, specifying pyxlsb as the engine
            df = pd.read_excel(file_data, engine="pyxlsb")
            # Convert the DataFrame to a text string
            text_rows = []
            for index, row in df.iterrows():
                text_rows.append(", ".join(row.astype(str)))
            text = "\n".join(text_rows)
        except Exception as e:
            self.logger.error(f"Error processing XLSB file at {link}: {e}")
            text = "Error processing XLSB file."

        return text

    def process_csv(self, link):
        try:
            import pandas as pd
        except ImportError:
            raise ImportError(
                "`pandas` package not found, please run `pip install pandas`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        try:
            # Assuming CSV uses default comma delimiter. If delimiter varies, consider detecting it.
            df = pd.read_csv(file_data, low_memory=False)
            # Convert the DataFrame to a text string, including headers
            text_rows = []
            for index, row in df.iterrows():
                text_rows.append(", ".join(row.astype(str)))
            text = "\n".join(text_rows)
        except Exception as e:
            self.logger.error(f"Error processing CSV file: {e}")
            text = "Error processing CSV file."

        return text

    def process_svg(self, link):
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            from reportlab.graphics import renderPM  # type: ignore
            from svglib.svglib import svg2rlg  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract`, `Pillow`, or `svglib` package not found, please run"
                " `pip install pytesseract Pillow svglib`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        drawing = svg2rlg(BytesIO(response.content))

        img_data = BytesIO()
        renderPM.drawToFile(drawing, img_data, fmt="PNG")
        img_data.seek(0)
        image = Image.open(img_data)

        return pytesseract.image_to_string(image)


if __name__ == "__main__":
    reader = ConfluenceReader()
