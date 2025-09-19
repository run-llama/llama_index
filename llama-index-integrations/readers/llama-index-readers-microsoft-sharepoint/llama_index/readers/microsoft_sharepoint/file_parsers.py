import logging
from typing import List, Union
from pathlib import Path

from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document

logger = logging.getLogger(__name__)


# PDF Reader
class PDFReader(BaseReader):
    """PDF reader using OCR for text extraction."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            import pytesseract
            from pdf2image import convert_from_path
        except ImportError:
            raise ImportError(
                "Please install pytesseract and pdf2image for PDFReader: pip install pytesseract pdf2image"
            )

        try:
            text = ""
            images = convert_from_path(str(file_path))
            for i, image in enumerate(images):
                image_text = pytesseract.image_to_string(image)
                text += f"Page {i + 1}:\n{image_text}\n\n"
            return [Document(text=text.strip(), metadata={"file_path": str(file_path)})]
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# HTML Reader
class HTMLReader(BaseReader):
    """HTML reader using BeautifulSoup for text extraction."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError(
                "Please install beautifulsoup4 for HTMLReader: pip install beautifulsoup4"
            )

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                html_content = f.read()
            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text(separator=" ", strip=True)
            return [Document(text=text, metadata={"file_path": str(file_path)})]
        except Exception as e:
            logger.error(f"Error processing HTML {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# TXT Reader
class TXTReader(BaseReader):
    """Plain text file reader."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            return [Document(text=text, metadata={"file_path": str(file_path)})]
        except Exception as e:
            logger.error(f"Error processing TXT {file_path}: {e}")
            # Try with different encoding
            try:
                with open(file_path, "r", encoding="latin-1") as f:
                    text = f.read()
                return [
                    Document(
                        text=text,
                        metadata={"file_path": str(file_path), "encoding": "latin-1"},
                    )
                ]
            except Exception as e2:
                logger.error(
                    f"Error processing TXT with fallback encoding {file_path}: {e2}"
                )
                return [
                    Document(
                        text="", metadata={"file_path": str(file_path), "error": str(e)}
                    )
                ]


# DOCX Reader
class DocxReader(BaseReader):
    """DOCX document reader."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            import docx2txt
        except ImportError:
            raise ImportError(
                "Please install docx2txt for DocxReader: pip install docx2txt"
            )

        try:
            text = docx2txt.process(str(file_path))
            return [Document(text=text or "", metadata={"file_path": str(file_path)})]
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# PPTX Reader
class PptxReader(BaseReader):
    """PowerPoint presentation reader."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "Please install python-pptx for PptxReader: pip install python-pptx"
            )

        try:
            text = ""
            presentation = Presentation(str(file_path))
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                text += slide_text + "\n"
            return [Document(text=text.strip(), metadata={"file_path": str(file_path)})]
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# CSV Reader
class CSVReader(BaseReader):
    """CSV file reader."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Please install pandas for CSVReader: pip install pandas")

        try:
            df = pd.read_csv(file_path, low_memory=False)
            # Include column headers
            text = f"Columns: {', '.join(df.columns.tolist())}\n\n"
            text_rows = []
            for _, row in df.iterrows():
                text_rows.append(", ".join(row.astype(str)))
            text += "\n".join(text_rows)
            return [
                Document(
                    text=text,
                    metadata={
                        "file_path": str(file_path),
                        "rows": len(df),
                        "columns": len(df.columns),
                    },
                )
            ]
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# XLSX Reader
class ExcelReader(BaseReader):
    """Excel spreadsheet reader."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError(
                "Please install pandas and openpyxl for ExcelReader: pip install pandas openpyxl"
            )

        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
            text = ""
            for sheet_name, sheet_data in sheets.items():
                text += f"Sheet: {sheet_name}\n"
                text += f"Columns: {', '.join(sheet_data.columns.tolist())}\n"
                for _, row in sheet_data.iterrows():
                    text += "\t".join(str(value) for value in row) + "\n"
                text += "\n"
            return [
                Document(
                    text=text.strip(),
                    metadata={"file_path": str(file_path), "sheets": len(sheets)},
                )
            ]
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# IMAGE Reader (OCR)
class ImageReader(BaseReader):
    """Image reader using OCR for text extraction."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError(
                "Please install pytesseract and Pillow for ImageReader: pip install pytesseract Pillow"
            )

        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return [
                Document(
                    text=text,
                    metadata={"file_path": str(file_path), "image_size": image.size},
                )
            ]
        except Exception as e:
            logger.error(f"Error processing Image {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# JSON Reader
class JSONReader(BaseReader):
    """JSON file reader."""

    def load_data(self, file_path: Union[str, Path], **kwargs) -> List[Document]:
        try:
            import json
        except ImportError:
            raise ImportError("JSON support should be built-in to Python")

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            # Convert JSON to readable text format
            text = json.dumps(data, indent=2, ensure_ascii=False)
            return [
                Document(
                    text=text, metadata={"file_path": str(file_path), "format": "json"}
                )
            ]
        except Exception as e:
            logger.error(f"Error processing JSON {file_path}: {e}")
            return [
                Document(
                    text="", metadata={"file_path": str(file_path), "error": str(e)}
                )
            ]


# Usage Example for SharePointReader:
# from .file_parsers import PDFReader, HTMLReader, DocxReader, PptxReader, CSVReader, ExcelReader, ImageReader, JSONReader, TXTReader
# custom_parsers = {
#     FileType.PDF: PDFReader(),
#     FileType.HTML: HTMLReader(),
#     FileType.DOCUMENT: DocxReader(),
#     FileType.PRESENTATION: PptxReader(),
#     FileType.CSV: CSVReader(),
#     FileType.SPREADSHEET: ExcelReader(),
#     FileType.IMAGE: ImageReader(),
#     FileType.JSON: JSONReader(),
#     FileType.TEXT: TXTReader(),
# }
# reader = SharePointReader(..., custom_parsers=custom_parsers, custom_folder="/tmp")
