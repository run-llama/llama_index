"""
Enhanced Slides parser.

Contains parsers for .pptx files with comprehensive content extraction.

"""

import gc
import logging
import os
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
from fsspec import AbstractFileSystem

from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.core.base.llms.base import BaseLLM
from llama_index.core import Settings

from .content_extractor import SlideContentExtractor

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)


class PptxReader(BaseReader):
    """
    Enhanced PowerPoint parser.

    Extract text, tables, charts, speaker notes, and optionally caption images.
    Supports multithreaded processing and LLM-based content consolidation.
    Always returns one Document per slide.
    """

    def __init__(
        self,
        extract_images: bool = False,
        context_consolidation_with_llm: bool = False,
        llm: Optional[BaseLLM] = None,
        batch_size: int = 10,
        num_workers: int = 4,
        raise_on_error: bool = False,
    ) -> None:
        """
        Initialize enhanced PptxReader.

        Args:
            extract_images: Whether to extract and caption images
            context_consolidation_with_llm: Whether to use LLM for contextual content consolidation
            llm: LLM instance for content consolidation (optional)
            batch_size: Number of slides to process in parallel batches
            num_workers: Number of worker threads (0 for sequential processing)

        """
        # Use provided LLM or fall back to global Settings
        if context_consolidation_with_llm and llm is None:
            llm = Settings.llm

        # Store settings
        self.extract_images = extract_images
        self.context_consolidation_with_llm = context_consolidation_with_llm
        self.llm = llm
        self.batch_size = batch_size
        self.num_workers = num_workers
        self.raise_on_error = raise_on_error
        self.content_extractor = SlideContentExtractor(
            llm=self.llm,
            extract_images=self.extract_images,
            context_consolidation_with_llm=self.context_consolidation_with_llm,
        )
        self.raise_on_error = raise_on_error

    def load_data(
        self,
        file: Union[str, Path],
        extra_info: Optional[Dict] = None,
        fs: Optional[AbstractFileSystem] = None,
    ) -> List[Document]:
        """
        Parse PowerPoint file with enhanced content extraction.

        Args:
            file: Path to the PowerPoint file
            extra_info: Additional metadata to include
            fs: File system to use for reading

        Returns:
            List of Documents (one per slide)

        """
        logger.debug(f"Loading PPTX file: {file}")
        file_path_str = str(file)

        # Extract content using enhanced extraction
        result = self.extract_with_validation(
            file_path=file_path_str,
            extract_images=self.extract_images,
            context_consolidation_with_llm=self.context_consolidation_with_llm,
            fs=fs,
        )

        if not result["success"] and not self.raise_on_error:
            logger.error(
                f"Failed to extract data from {file_path_str}: {result['errors']}"
            )
            return []
        elif not result["success"] and self.raise_on_error:
            raise ValueError(
                f"Failed to extract data from {file_path_str}: {result['errors']}"
            )

        # Convert to Documents
        docs = []
        for i, slide in enumerate(result["data"]["slides"], start=1):
            # Create rich metadata
            metadata = {
                "file_path": str(file),
                "page_label": i,
                "title": slide.get("title", ""),
                "extraction_errors": slide.get("extraction_errors", []),
                "extraction_warnings": slide.get("extraction_warnings", []),
                "tables": slide.get("tables", []),
                "charts": slide.get("charts", []),
                "notes": slide.get("notes", ""),
                "images": slide.get("images", []),
                "text_sections": slide.get("text_sections", []),
            }
            if extra_info:
                metadata.update(extra_info)

            docs.append(
                Document(
                    text=slide["content"],
                    metadata=metadata,
                    excluded_embed_metadata_keys=list(
                        metadata.keys()
                    ),  # excluding the metadata keys from the embedding since the metadata size can potentially be too large and may cause the embedding to fail
                    excluded_llm_metadata_keys=list(
                        metadata.keys()
                    ),  # excluding the metadata keys from the llm
                )
            )

        logger.debug(f"Successfully loaded {len(docs)} slides from {file}")
        return docs

    def extract_with_validation(
        self,
        file_path: str,
        extract_images: bool = True,
        context_consolidation_with_llm: bool = False,
        fs: Optional[AbstractFileSystem] = None,
    ) -> Dict[str, Any]:
        """Extract content from PowerPoint file with validation and multithreaded processing."""
        result: Dict[str, Any] = {
            "success": False,
            "data": None,
            "errors": [],
            "warnings": [],
            "stats": {},
        }

        # Validate file and get presentation object
        validation = self._validate_file(file_path, fs)
        if not validation.get("valid", False):
            result["errors"] = validation.get("errors", [])
            return result

        # Use the presentation object from validation
        presentation = validation.get("presentation")
        if presentation is None:
            result["errors"].append("Failed to get presentation object from validation")
            return result

        filename = Path(file_path).name
        logger.debug(f"Processing file: {filename}")

        try:
            start_time = datetime.now()
            total_slides = len(presentation.slides)
            logger.debug(f"Processing {total_slides} slides from {filename}")

            # Prepare result structure
            slides_data: list = []

            # Create batches of slide indices
            batches = [
                (i, min(i + self.batch_size, total_slides))
                for i in range(0, total_slides, self.batch_size)
            ]

            # Process in parallel or serial
            if self.num_workers and self.num_workers > 0:
                with ThreadPoolExecutor(max_workers=self.num_workers) as executor:
                    futures = {
                        executor.submit(
                            self._process_batch,
                            presentation,
                            start,
                            end,
                            filename,
                            extract_images,
                            context_consolidation_with_llm,
                        ): (start, end)
                        for start, end in batches
                    }
                    for fut in as_completed(futures):
                        start, end = futures[fut]
                        try:
                            batch_results = fut.result()
                            slides_data.extend(batch_results)
                        except Exception as e:
                            logger.error(f"Batch {start + 1}-{end} failed: {e}")
                            for idx in range(start, end):
                                slides_data.append(
                                    {
                                        "slide_number": idx + 1,
                                        "error": str(e),
                                        "partial_extraction": True,
                                    }
                                )
                        finally:
                            gc.collect()
            else:
                # Serial fallback
                for start, end in batches:
                    slides_data.extend(
                        self._process_batch(
                            presentation,
                            start,
                            end,
                            filename,
                            extract_images,
                            context_consolidation_with_llm,
                        )
                    )
                    gc.collect()

            # Calculate stats and finalize
            processing_time = (datetime.now() - start_time).total_seconds()
            stats = {
                "total_slides": total_slides,
                "processed_slides": len(slides_data),
                "total_errors": sum(
                    1
                    for s in slides_data
                    if s.get("error") or s.get("extraction_errors")
                ),
                "processing_time_seconds": processing_time,
                "file_size_mb": 0 if fs else os.path.getsize(file_path) / (1024 * 1024),
            }

            result.update(
                {
                    "success": True,
                    "data": {
                        "filename": filename,
                        "slides": sorted(
                            slides_data, key=lambda s: s.get("slide_number", 0)
                        ),
                        "metadata": {
                            "total_slides": total_slides,
                            "file_path": file_path,
                            "processing_timestamp": datetime.now().isoformat(),
                            "extract_images": extract_images,
                            "context_consolidation_with_llm": context_consolidation_with_llm,
                        },
                    },
                    "errors": [],
                    "stats": stats,
                }
            )

            logger.debug(
                f"Successfully processed {filename}: {stats['processed_slides']} slides "
                f"in {stats['processing_time_seconds']:.2f}s"
            )

            # Cleanup
            del presentation
            gc.collect()

        except Exception as e:
            error_msg = f"Critical error processing {filename}: {e}"
            logger.error(error_msg)
            result["errors"].append(error_msg)

        return result

    def _process_batch(
        self,
        presentation: Any,
        start: int,
        end: int,
        filename: str,
        extract_images: bool,
        context_consolidation_with_llm: bool,
    ) -> list:
        """
        Process slides in the range [start, end) and return their extracted data.
        Runs in the context of a worker thread.
        """
        thread_name = threading.current_thread().name
        logger.debug(f"[{thread_name}] Starting batch {start + 1}-{end}")

        batch_data: list = []
        for idx in range(start, end):
            try:
                slide = presentation.slides[idx]
                slide_data = self.content_extractor.extract_slide_safe(
                    slide=slide,
                    slide_number=idx + 1,
                    filename=Path(filename),
                )
                batch_data.append(slide_data)
                if slide_data.get("extraction_errors"):
                    logger.warning(
                        f"[{thread_name}] Slide {idx + 1} had extraction errors: {slide_data.get('extraction_errors')}"
                    )
            except Exception as e:
                logger.warning(f"[{thread_name}] Error on slide {idx + 1}: {e}")
                batch_data.append(
                    {
                        "slide_number": idx + 1,
                        "error": str(e),
                        "partial_extraction": True,
                    }
                )
        logger.debug(f"[{thread_name}] Finished batch {start + 1}-{end}")
        return batch_data

    def _validate_file(
        self, file_path: str, fs: Optional[AbstractFileSystem] = None
    ) -> Dict[str, Any]:
        """Validate that the file exists, and can be opened. Returns presentation object for reuse."""
        from pptx import Presentation
        import io

        validation: Dict[str, Any] = {
            "valid": True,
            "errors": [],
            "warnings": [],
            "presentation": None,
        }

        # Extension warning
        if not file_path.lower().endswith((".pptx", ".ppt")):
            validation["warnings"].append(
                "File extension not typical for PowerPoint files"
            )

        # File existence (only for local files)
        if not fs:
            if not os.path.exists(file_path):
                validation["valid"] = False
                validation["errors"].append(f"File not found: {file_path}")
                return validation

        # Try opening the presentation
        try:
            if fs:
                with fs.open(file_path) as f:
                    presentation = Presentation(io.BytesIO(f.read()))
            else:
                presentation = Presentation(file_path)

            count = len(presentation.slides)
            if count == 0:
                validation["warnings"].append("Presentation contains no slides")
            elif count > 1000:
                validation["warnings"].append(f"Large presentation: {count} slides")

            # Return the presentation object for reuse
            validation["presentation"] = presentation

        except Exception as e:
            validation["valid"] = False
            validation["errors"].append(f"Cannot open as PowerPoint file: {e}")

        return validation
