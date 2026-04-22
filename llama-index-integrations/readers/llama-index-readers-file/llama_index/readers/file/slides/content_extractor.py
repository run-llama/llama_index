"""
Enhanced content extraction for PowerPoint slides.

Extracts text, tables, charts, and speaker notes with proper structure preservation.
"""

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from tenacity import (
    retry,
    retry_if_exception,
    wait_exponential,
    stop_after_attempt,
    before_sleep_log,
)

from llama_index.core.base.llms.base import BaseLLM

from .image_extractor import ImageExtractor

logger = logging.getLogger(__name__)

# Retry constants
MAX_ATTEMPTS = 5
MIN_BACKOFF = 1  # seconds
MAX_BACKOFF = 60  # seconds

# Content formatting constants
SECTION_SEPARATOR = "-----"


def is_rate_limit_error(exception: Exception) -> bool:
    """Check if the exception is a rate limit error."""
    error_str = str(exception).lower()
    return any(
        phrase in error_str
        for phrase in [
            "rate limit",
            "too many requests",
            "429",
            "quota exceeded",
            "rate_limit_exceeded",
            "throttle",
            "rate limiting",
        ]
    )


class SlideContentExtractor:
    """
    Enhanced content extractor for PowerPoint slides.

    Extracts comprehensive content including tables, charts, notes, and structured text.
    """

    def __init__(
        self,
        llm: Optional[BaseLLM] = None,
        extract_images: bool = True,
        context_consolidation_with_llm: bool = False,
    ):
        """
        Initialize content extractor.

        Args:
            llm: Optional LLM for content consolidation
            extract_images: Whether to extract and caption images
            context_consolidation_with_llm: Whether to use LLM for content consolidation

        """
        self.llm = llm
        self.extract_images = extract_images
        self.context_consolidation_with_llm = context_consolidation_with_llm

        # Initialize image extractor if needed
        self.image_extractor = None
        if self.extract_images:
            self.image_extractor = ImageExtractor()

    def extract_slide_safe(
        self,
        slide,
        slide_number: int,
        filename: Path,
    ) -> Dict[str, Any]:
        """Extract slide content with comprehensive error recovery."""
        slide_data = {
            "slide_number": slide_number,
            "title": "",
            "text_sections": [],
            "tables": [],
            "charts": [],
            "notes": "",
            "images": [],
            "extraction_errors": [],
            "extraction_warnings": [],
            "filename": str(filename),
        }

        # Extract title safely
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_data["title"] = slide.shapes.title.text.strip()
            else:
                # Fallback: find likely title from slide content
                slide_data["title"] = self._find_likely_title(slide)
        except Exception:
            # Final fallback: find likely title from slide content
            slide_data["title"] = self._find_likely_title(slide)

        # Process all shapes with error isolation
        for shape_idx, shape in enumerate(slide.shapes):
            self._process_shape_safe(shape, slide_data)

        # Extract speaker notes
        try:
            slide_data["notes"] = self._extract_slide_notes(slide)
        except Exception as e:
            # Notes are supplementary - add warning instead of error
            slide_data["extraction_warnings"].append(
                {"component": "notes", "warning": f"Notes extraction failed: {e!s}"}
            )
            slide_data["notes"] = ""  # Set to empty string
        try:
            slide_data["content"] = self._consolidate_slide_data(
                slide_data, with_llm=self.context_consolidation_with_llm
            )
        except Exception as e:
            slide_data["extraction_errors"].append(
                {
                    "component": "content_consolidation",
                    "error": f"Content consolidation failed: {e!s}",
                }
            )
            # Fallback to basic consolidation
            slide_data["content"] = self._consolidate_slide_data(
                slide_data, with_llm=False
            )

        return slide_data

    def _find_likely_title(self, slide) -> str:
        """Find the most likely title from slide shapes using balanced scoring."""
        try:
            candidates = []

            # Get slide height from presentation (slide.part.presentation gives us the presentation)
            try:
                presentation = slide.part.presentation
                slide_height = presentation.slide_height
            except Exception:
                slide_height = 6858000  # Default slide height in EMUs

            for shape in slide.shapes:
                if (
                    hasattr(shape, "text_frame")
                    and shape.text_frame
                    and shape.text_frame.text.strip()
                ):
                    text = shape.text_frame.text.strip()
                    score = 0

                    # Position score (0-100): closer to top = higher score
                    # The value of top is basically the distance from the top of the slide to the top of the shape, so smaller top = closer to top of the slide
                    if hasattr(shape, "top") and shape.top is not None:
                        position_score = max(
                            0, 100 * (slide_height - shape.top) / slide_height
                        )
                        score += position_score

                    # Size score (0-50): smaller height = higher score (titles are compact)
                    # The value of height is basically the height of the shape, so smaller height = higher score, as title would ideally be compact
                    if hasattr(shape, "height") and shape.height is not None:
                        size_score = max(0, 50 * (1 - shape.height / slide_height))
                        score += size_score

                    # Text length score (0-25): shorter text = higher score, as titles are typically short
                    if len(text) < 100:
                        text_score = min(25, 25 - len(text) / 4)
                        score += text_score

                    candidates.append((text, score))

            if candidates:
                # Return highest scoring candidate
                candidates.sort(key=lambda x: x[1], reverse=True)
                return candidates[0][0]

            return ""
        except Exception:
            return ""

    def _process_shape_safe(self, shape, slide_data: Dict[str, Any]) -> None:
        """Safely process shape with error handling."""
        try:
            self._process_shape(shape, slide_data)
        except Exception as e:
            # Log error but continue processing
            slide_data["extraction_errors"].append(
                {
                    "shape_processing": str(e),
                    "shape_type": getattr(shape, "shape_type", "unknown"),
                }
            )

    def _process_shape(self, shape, slide_data: Dict[str, Any]) -> None:
        """Process individual shape content extraction."""
        # Extract text content
        if hasattr(shape, "text") and shape.text.strip():
            text_data = self._extract_text_hierarchical(shape)
            if text_data:
                slide_data["text_sections"].append(text_data)

        # Extract table content
        if shape.has_table:
            table_data = self._extract_table_content(shape.table)
            slide_data["tables"].append(table_data)

        # Extract chart content
        if shape.has_chart:
            chart_data = self._extract_chart_metadata(shape.chart)
            slide_data["charts"].append(chart_data)

        # Extract image content if enabled
        if self.extract_images and self.image_extractor:
            if hasattr(shape, "image") and shape.image:
                try:
                    image_data = self.image_extractor.extract_image_data(
                        shape, slide_data["slide_number"]
                    )
                    slide_data["images"].append(image_data)
                except Exception as e:
                    error_msg = f"Image extraction failed: {e!s}"
                    logger.debug(error_msg)
                    slide_data["extraction_warnings"].append(error_msg)

        # Process grouped shapes recursively
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shape_in_group in shape.shapes:
                self._process_shape_safe(shape_in_group, slide_data)

    def _extract_text_hierarchical(self, shape) -> Dict[str, Any]:
        """Extract text preserving structure and formatting context."""
        if not shape.has_text_frame:
            return None

        text_data = {
            "content": "",
            "structure": [],
            "enhanced_content": "",
            "formatting_preserved": True,
        }

        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            para_info = {
                "text": "",
                "level": paragraph.level,  # Bullet/indent level (0-8)
                "runs": [],
            }

            # Build enhanced text for this paragraph
            enhanced_para_text = ""

            # Add indentation based on level
            indent = "  " * paragraph.level  # 2 spaces per level

            for run in paragraph.runs:
                run_text = run.text
                para_info["text"] += run_text

                # Preserve formatting metadata for enhanced context
                run_info = {
                    "text": run_text,
                    "bold": bool(run.font.bold),
                    "italic": bool(run.font.italic),
                    "size": run.font.size.pt if run.font.size else None,
                }
                para_info["runs"].append(run_info)

                # Create enhanced text with formatting cues
                enhanced_run_text = run_text
                if run.font.bold:
                    enhanced_run_text = f"**{enhanced_run_text}**"
                if run.font.italic:
                    enhanced_run_text = f"*{enhanced_run_text}*"
                if (
                    run.font.size and run.font.size.pt > 14
                ):  # Larger text (likely headings)
                    enhanced_run_text = f"# {enhanced_run_text}"

                enhanced_para_text += enhanced_run_text

            text_data["structure"].append(para_info)
            text_data["content"] += para_info["text"] + "\n"

            # Add enhanced paragraph with indentation and formatting
            if enhanced_para_text.strip():
                text_data["enhanced_content"] += (
                    f"{indent}{enhanced_para_text.strip()}\n"
                )

        return text_data

    def _extract_table_content(self, table) -> Dict[str, Any]:
        """Extract table preserving structure for semantic search."""
        table_data = {
            "type": "table",
            "dimensions": {"rows": len(table.rows), "columns": len(table.columns)},
            "headers": [],
            "data": [],
            "detailed_content": "",
            "narrative_description": "",
        }

        # Identify and extract headers
        has_header_row = getattr(table, "first_row", False)
        header_row_idx = 0 if has_header_row else -1

        if has_header_row and table.rows:
            for cell in table.rows[0].cells:
                header_text = self._extract_cell_text(cell)
                table_data["headers"].append(header_text)

        # Extract data with context preservation
        for row_idx, row in enumerate(table.rows):
            if row_idx == header_row_idx:
                continue

            row_data = []
            row_text_parts = []

            for col_idx, cell in enumerate(row.cells):
                cell_text = self._extract_cell_text(cell)
                row_data.append(cell_text)

                # Create contextual text for this cell
                if table_data["headers"] and col_idx < len(table_data["headers"]):
                    header = table_data["headers"][col_idx]
                    row_text_parts.append(f"{header}: {cell_text}")
                else:
                    row_text_parts.append(cell_text)

            table_data["data"].append(row_data)

            # Build detailed content with context
            row_context = " | ".join(row_text_parts)
            table_data["detailed_content"] += row_context + "\n"

        # Create narrative description for better comprehension
        table_data["narrative_description"] = self._create_table_narrative(table_data)

        return table_data

    def _extract_cell_text(self, cell) -> str:
        """Extract text from a table cell."""
        try:
            text_frame = cell.text_frame
            if text_frame is None:
                return ""
            return text_frame.text.strip()
        except Exception:
            return ""

    def _create_table_narrative(self, table_data: Dict) -> str:
        """Convert table to narrative form for better RAG comprehension."""
        narrative_parts = []

        if table_data["headers"]:
            headers_text = ", ".join(table_data["headers"])
            narrative_parts.append(f"Table with columns: {headers_text}")

        narrative_parts.append(f"Contains {len(table_data['data'])} data rows")

        # Add all table data
        if table_data["data"]:
            for row_idx, sample_row in enumerate(table_data["data"]):
                if table_data["headers"]:
                    sample_parts = []
                    for col_idx, value in enumerate(
                        sample_row[: len(table_data["headers"])]
                    ):
                        if col_idx < len(table_data["headers"]):
                            sample_parts.append(
                                f"{table_data['headers'][col_idx]}={value}"
                            )
                    narrative_parts.append(
                        f"Row {row_idx + 1}: {', '.join(sample_parts)}"
                    )
                else:
                    # If no headers, just show the row values
                    narrative_parts.append(
                        f"Row {row_idx + 1}: {', '.join(sample_row)}"
                    )

        return ". ".join(narrative_parts)

    def _extract_chart_metadata(self, chart) -> Dict[str, Any]:
        """Extract available chart metadata with fallback strategies, including series values."""
        chart_data: Dict[str, Any] = {
            "type": "chart",
            "chart_type": (
                str(chart.chart_type) if hasattr(chart, "chart_type") else "unknown"
            ),
            "title": "",
            "axes": {},
            "series_info": [],
            "description": "",
            "extraction_limitations": [],
            "content": "",
        }

        # Extract chart title
        try:
            if chart.has_title and chart.chart_title:
                chart_data["title"] = chart.chart_title.text_frame.text
        except Exception:
            chart_data["extraction_limitations"].append("Title extraction failed")

        # Extract axis information
        try:
            if (
                hasattr(chart, "category_axis")
                and chart.category_axis
                and chart.category_axis.has_title
            ):
                chart_data["axes"]["category"] = {
                    "title": chart.category_axis.axis_title.text_frame.text,
                    "visible": chart.category_axis.visible,
                }
            if (
                hasattr(chart, "value_axis")
                and chart.value_axis
                and chart.value_axis.has_title
            ):
                chart_data["axes"]["value"] = {
                    "title": chart.value_axis.axis_title.text_frame.text,
                    "visible": chart.value_axis.visible,
                }
        except Exception:
            chart_data["extraction_limitations"].append("Axis extraction limited")

        # Extract series, categories, and values
        try:
            for plot in chart.plots:
                # get categories
                try:
                    categories: List[str] = [str(cat) for cat in plot.categories]
                except Exception:
                    categories = []
                    chart_data["extraction_limitations"].append(
                        "Category extraction failed"
                    )

                # get each series' name and values
                series_list: List[Dict[str, Any]] = []
                for series in plot.series:
                    name = str(series.name) if series.name else ""
                    try:
                        values = [float(v) for v in series.values]
                    except Exception:
                        values = []
                        chart_data["extraction_limitations"].append(
                            f'Values extraction failed for series "{name}"'
                        )
                    series_list.append({"name": name, "values": values})

                chart_data["series_info"].append(
                    {
                        "series_count": len(series_list),
                        "categories": categories,
                        "series": series_list,
                    }
                )
        except Exception:
            chart_data["extraction_limitations"].append(
                "Series data extraction limited"
            )

        # Generate descriptive text for RAG
        chart_data["description"] = self._generate_chart_description(chart_data)
        chart_data["content"] = chart_data[
            "description"
        ]  # For consistency with other content types

        return chart_data

    def _generate_chart_description(self, chart_data: Dict[str, Any]) -> str:
        """Generate a comprehensive description of the chart for RAG purposes."""
        description_parts = []

        # Chart type and title
        chart_type = chart_data.get("chart_type", "unknown")
        title = chart_data.get("title", "")
        if title:
            description_parts.append(f"{chart_type} chart titled '{title}'")
        else:
            description_parts.append(f"{chart_type} chart")

        # Axes information
        axes = chart_data.get("axes", {})
        if "category" in axes:
            description_parts.append(f"Category axis: {axes['category']['title']}")
        if "value" in axes:
            description_parts.append(f"Value axis: {axes['value']['title']}")

        # Series data
        series_info = chart_data.get("series_info", [])
        for plot_idx, plot_info in enumerate(series_info):
            categories = plot_info.get("categories", [])
            series_list = plot_info.get("series", [])

            for series in series_list:
                name = series.get("name", "")
                values = series.get("values", [])
                if name and values:
                    # Pair categories with values if available
                    if categories and len(categories) >= len(values):
                        pairs = ", ".join(
                            f"{cat}={val}" for cat, val in zip(categories, values)
                        )
                        description_parts.append(f"Series '{name}': {pairs}")
                    else:
                        # Fallback to just values if no categories or mismatch
                        value_preview = ", ".join([str(v) for v in values])
                        description_parts.append(f"Series '{name}': {value_preview}")

                    # Add summary stats
                    try:
                        min_val = min(values)
                        max_val = max(values)
                        avg_val = sum(values) / len(values)
                        description_parts.append(
                            f"Series '{name}' Range: {min_val:.2f} to {max_val:.2f}, Average: {avg_val:.2f}"
                        )
                    except (ValueError, ZeroDivisionError):
                        pass

        # Extraction limitations
        limitations = chart_data.get("extraction_limitations", [])
        if limitations:
            description_parts.append(f"Note: {', '.join(limitations)}")

        return ". ".join(description_parts)

    def _extract_slide_notes(self, slide) -> str:
        """Extract speaker notes with proper error handling."""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()

            # Add contextual prefix
            if notes_text:
                return f"[Speaker Notes]: {notes_text}"
        return ""

    def _consolidate_slide_data(
        self, slide_data: Dict[str, Any], with_llm: bool = False
    ) -> str:
        """Consolidate slide data into a single text block"""
        consolidated_text = []

        if slide_data["title"]:
            consolidated_text.append(f"Title: {slide_data['title']}")

        # Add text sections - use enhanced_content if available for better formatting context
        text_sections = []
        for section in slide_data["text_sections"]:
            if section.get("enhanced_content"):
                text_sections.append(section["enhanced_content"].strip())
            else:
                text_sections.append(section["content"])

        if text_sections:
            if consolidated_text:  # Add separator if we already have content
                consolidated_text.append(SECTION_SEPARATOR)
            consolidated_text.extend(text_sections)

        # Add table descriptions
        table_sections = []
        for table in slide_data["tables"]:
            table_sections.append(f"Table: {table.get('narrative_description', '')}")

        if table_sections:
            if consolidated_text:
                consolidated_text.append(SECTION_SEPARATOR)
            consolidated_text.extend(table_sections)

        # Add chart descriptions
        chart_sections = []
        for chart in slide_data["charts"]:
            chart_sections.append(f"Chart: {chart.get('description', '')}")

        if chart_sections:
            if consolidated_text:
                consolidated_text.append(SECTION_SEPARATOR)
            consolidated_text.extend(chart_sections)

        # Add image descriptions
        image_sections = []
        for image in slide_data["images"]:
            if "error" not in image:
                image_sections.append(
                    f"Image: {image.get('caption', 'No caption available')}"
                )

        if image_sections:
            if consolidated_text:
                consolidated_text.append(SECTION_SEPARATOR)
            consolidated_text.extend(image_sections)

        # Add speaker notes
        if slide_data["notes"]:
            if consolidated_text:
                consolidated_text.append(SECTION_SEPARATOR)
            consolidated_text.append(slide_data["notes"])

        slide_content = "\n".join(consolidated_text).strip()

        if slide_content and with_llm:
            prompt = """
            Act as a senior business analyst preparing notes for this presentation slide. Your goal is to produce a single, comprehensive paragraph that captures the entire slide's content and its intended strategic message.

            To achieve this, synthesize all the provided information by following these rules:

            1.  **Be Methodically Comprehensive:** Systematically describe all content on the slide, including its title, all text points, data from any tables, and key details from charts. Ensure no information is omitted.

            2.  **Synthesize into a Unified Narrative:** Weave all the extracted details into a single, flowing paragraph. Your summary must read like a cohesive story that explains the slide, not like a disconnected list of facts.

            3.  **Identify the Core Conclusion:** Clearly state the main takeaway or the specific conclusion the presenter wants the audience to draw from the data presented.

            4.  **Adhere Strictly to the Source:** Base your summary *only* on the information provided in the slide. Do not make assumptions or introduce any external knowledge.

            5.  **Deliver as Plain Text:** The final output must be one continuous paragraph of plain English text, without any special formatting like backticks, bullet points, or section headers.
            The paragraph should be concise, accurate, and in English.

            6. The output should not exceed 300 words.
            """
            try:
                refined_slide_content = self._chat(
                    [
                        prompt
                        + "\n\n"
                        + "<Slide Content>"
                        + slide_content
                        + "\n\n</Slide Content>"
                    ]
                )
            except Exception as e:
                return slide_content

            return (
                refined_slide_content.strip()
                if refined_slide_content
                else slide_content
            )

        return slide_content if slide_content else ""

    @retry(
        retry=retry_if_exception(is_rate_limit_error),
        wait=wait_exponential(multiplier=MIN_BACKOFF, min=MIN_BACKOFF, max=MAX_BACKOFF),
        stop=stop_after_attempt(MAX_ATTEMPTS),
        before_sleep=before_sleep_log(logger, logging.WARNING),
    )
    def _chat(self, messages: List[str]) -> str:
        """Use LLM to process messages and return response with retry logic for rate limits."""
        if not self.llm:
            raise RuntimeError("LLM not available for chat processing")

        # Combine messages into a single prompt
        full_prompt = "\n".join(messages)
        response = self.llm.complete(full_prompt)
        return response.text
