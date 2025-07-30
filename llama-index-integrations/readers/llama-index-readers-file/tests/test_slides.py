import pytest
from llama_index.readers.file.slides import PptxReader
from .generate_test_ppt import create_comprehensive_test_presentation


@pytest.fixture()
def pptx_file(tmp_path):
    """Create a temporary PowerPoint file for testing."""
    if create_comprehensive_test_presentation is None:
        pytest.skip("generate_test_ppt not available")

    # Create test presentation in temp directory
    file_path = tmp_path / "test_presentation.pptx"
    create_comprehensive_test_presentation(str(file_path))
    return file_path


def test_pptx_reader_init():
    """Test PptxReader initialization."""
    reader = PptxReader(extract_images=False, num_workers=2)
    assert reader.extract_images is False
    assert reader.num_workers == 2


def test_load_data_pptx(pptx_file):
    """Test loading PowerPoint data."""
    reader = PptxReader(extract_images=False, context_consolidation_with_llm=False)

    documents = reader.load_data(pptx_file)

    # Basic validation
    assert len(documents) == 12  # Should have 12 slides
    assert all(hasattr(doc, "text") for doc in documents)
    assert all(hasattr(doc, "metadata") for doc in documents)

    # Check first slide
    first_slide = documents[0]
    assert "Enhanced PowerPoint Reader Test" in first_slide.text
    assert first_slide.metadata["page_label"] == 1
    assert len(first_slide.metadata["notes"]) > 0  # Check notes content exists


def test_table_extraction(pptx_file):
    """Test table data extraction."""
    reader = PptxReader()
    documents = reader.load_data(pptx_file)

    # Find slides with tables
    table_slides = [doc for doc in documents if len(doc.metadata.get("tables", [])) > 0]
    assert len(table_slides) >= 1

    # Check table metadata
    table_slide = table_slides[0]
    tables = table_slide.metadata.get("tables", [])
    assert len(tables) > 0

    table = tables[0]
    assert "headers" in table
    assert "data" in table
    assert "dimensions" in table


def test_chart_extraction(pptx_file):
    """Test chart metadata extraction."""
    reader = PptxReader()
    documents = reader.load_data(pptx_file)

    # Find slides with charts
    chart_slides = [doc for doc in documents if len(doc.metadata.get("charts", [])) > 0]
    assert len(chart_slides) >= 1

    # Check chart metadata
    chart_slide = chart_slides[0]
    charts = chart_slide.metadata.get("charts", [])
    assert len(charts) > 0

    chart = charts[0]
    assert "chart_type" in chart
    assert "series_info" in chart


def test_speaker_notes_extraction(pptx_file):
    """Test speaker notes extraction."""
    reader = PptxReader()
    documents = reader.load_data(pptx_file)

    # All slides should have notes
    slides_with_notes = [
        doc for doc in documents if len(doc.metadata.get("notes", "")) > 0
    ]
    assert len(slides_with_notes) == 12

    # Check notes content
    first_slide = documents[0]
    notes = first_slide.metadata.get("notes", "")
    assert len(notes) > 0
    assert "comprehensive test presentation" in notes.lower()


def test_content_consolidation(pptx_file):
    """Test content consolidation structure."""
    reader = PptxReader()
    documents = reader.load_data(pptx_file)

    # Check content structure
    for doc in documents:
        assert "-----" in doc.text  # Section separators
        assert len(doc.text) > 0  # Content should exist


def test_multithreading(pptx_file):
    """Test multithreaded processing."""
    reader = PptxReader(num_workers=2, batch_size=4)
    documents = reader.load_data(pptx_file)

    # Should process successfully with threading
    assert len(documents) == 12
    assert all(doc.metadata.get("page_label") for doc in documents)


def test_llm_consolidation_with_settings_llm(pptx_file):
    """Test LLM consolidation when LLM is set in Settings but not passed directly."""
    from llama_index.core import Settings
    from llama_index.core.llms.mock import MockLLM

    Settings.llm = MockLLM()
    reader = PptxReader(
        extract_images=False,
        context_consolidation_with_llm=True,  # Request LLM consolidation
        llm=None,  # Don't pass LLM directly
        num_workers=2,
    )

    # Should still return results
    documents = reader.load_data(pptx_file)

    # Basic validation
    assert len(documents) == 12
    assert all(hasattr(doc, "text") for doc in documents)


def test_llm_consolidation_with_direct_llm(pptx_file):
    """Test LLM consolidation when LLM is passed directly to PptxReader."""
    from llama_index.core.llms.mock import MockLLM

    # Create MockLLM directly and pass it to the reader
    mock_llm = MockLLM()

    reader = PptxReader(
        extract_images=False,
        context_consolidation_with_llm=True,  # Request LLM consolidation
        llm=mock_llm,  # Pass LLM directly
        num_workers=2,
    )

    # Should use the directly passed LLM
    assert reader.context_consolidation_with_llm is True
    assert reader.llm is mock_llm  # Should be the exact same instance

    # Should process successfully
    documents = reader.load_data(pptx_file)

    # Basic validation
    assert len(documents) == 12
    assert all(hasattr(doc, "text") for doc in documents)
    assert all(hasattr(doc, "metadata") for doc in documents)

    # Content should be consolidated
    for doc in documents:
        assert "-----" in doc.text  # Section separators should be there
        assert len(doc.text) > 0  # Should have content


def test_title_detection_scoring(tmp_path):
    """Test that title detection correctly identifies titles using position and size scoring."""
    from pptx import Presentation
    from pptx.util import Inches

    # Create a test presentation with different text positions
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title at top (should win)
    title_box = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(0.8)
    )
    title_box.text_frame.text = "Test Title"

    # Add body text in middle (larger box, should lose)
    body_box = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(2)
    )
    body_box.text_frame.text = (
        "This is longer body text that should not be detected as title"
    )

    # Add footer at bottom (should lose due to position)
    footer_box = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(6), width=Inches(8), height=Inches(0.5)
    )
    footer_box.text_frame.text = "Footer"

    # Save test presentation
    test_file = tmp_path / "title_test.pptx"
    prs.save(str(test_file))

    # Test title detection
    reader = PptxReader(extract_images=False, context_consolidation_with_llm=False)
    documents = reader.load_data(test_file)

    # Should have one document
    assert len(documents) == 1
    doc = documents[0]
    # Title should be detected correctly
    assert doc.metadata["title"] == "Test Title"


def test_title_detection_edge_cases(tmp_path):
    """Test title detection edge cases through the public API."""
    from pptx import Presentation
    from pptx.util import Inches

    # Test 1: Empty slide (no shapes)
    prs = Presentation()
    empty_slide = prs.slides.add_slide(prs.slide_layouts[6])

    test_file = tmp_path / "empty_test.pptx"
    prs.save(str(test_file))

    reader = PptxReader(extract_images=False, context_consolidation_with_llm=False)
    documents = reader.load_data(test_file)

    assert len(documents) == 1
    assert documents[0].metadata["title"] == ""  # Empty slide should have empty title

    # Test 2: Multiple titles at different positions (top should win)
    prs2 = Presentation()
    multi_title_slide = prs2.slides.add_slide(prs2.slide_layouts[6])

    # Title 1 at very top
    top_title = multi_title_slide.shapes.add_textbox(
        left=Inches(1), top=Inches(0.2), width=Inches(8), height=Inches(0.6)
    )
    top_title.text_frame.text = "Top Title"

    # Title 2 lower but smaller
    lower_title = multi_title_slide.shapes.add_textbox(
        left=Inches(1), top=Inches(1.5), width=Inches(6), height=Inches(0.4)
    )
    lower_title.text_frame.text = "Lower"

    test_file2 = tmp_path / "multi_title_test.pptx"
    prs2.save(str(test_file2))

    documents2 = reader.load_data(test_file2)
    assert len(documents2) == 1
    assert documents2[0].metadata["title"] == "Top Title"  # Top position should win


def test_raise_on_error_parameter(tmp_path):
    """Test raise_on_error parameter behavior with invalid files."""
    # Test 1: raise_on_error=False (default) - should return empty list on error
    reader_no_raise = PptxReader(extract_images=False, raise_on_error=False)

    # Try to read a non-existent file
    non_existent_file = tmp_path / "does_not_exist.pptx"
    documents = reader_no_raise.load_data(non_existent_file)
    assert documents == []  # Should return empty list, not raise error

    # Test 2: raise_on_error=True - should raise ValueError on error
    reader_with_raise = PptxReader(extract_images=False, raise_on_error=True)

    # Try to read the same non-existent file, should raise ValueError
    with pytest.raises(ValueError, match="Failed to extract data"):
        reader_with_raise.load_data(non_existent_file)

    # Test 3: Create an invalid file and test both behaviors
    invalid_file = tmp_path / "invalid.pptx"
    invalid_file.write_text("This is not a valid PowerPoint file")

    # With raise_on_error=False, should return empty list
    documents = reader_no_raise.load_data(invalid_file)
    assert documents == []

    # With raise_on_error=True, should raise ValueError
    with pytest.raises(ValueError, match="Failed to extract data"):
        reader_with_raise.load_data(invalid_file)
