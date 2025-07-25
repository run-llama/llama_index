"""
Generate a comprehensive PowerPoint presentation for testing the enhanced PptxReader.

This script creates a presentation with:
- Text content with formatting
- Tables with structured data
- Charts with series data
- Speaker notes
- Mixed content slides

Perfect for demonstrating the enhanced extraction capabilities.
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_comprehensive_test_presentation(
    filename="comprehensive_test_presentation.pptx",
):
    """Create a comprehensive test PowerPoint presentation."""
    # Create presentation
    prs = Presentation()

    # Slide 1: Title slide with rich text
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Enhanced PowerPoint Reader Test"
    subtitle.text = (
        "Comprehensive Content Extraction Demo\nTesting Tables, Charts, Notes & Text"
    )

    # Add speaker notes
    notes_slide = slide1.notes_slide
    notes_slide.notes_text_frame.text = (
        "Welcome to the comprehensive test presentation. "
        "This presentation demonstrates the enhanced PowerPoint reader's ability to extract "
        "various content types including formatted text, tables, charts, and speaker notes. "
        "Each slide showcases different content extraction scenarios."
    )

    # Slide 2: Rich text content with formatting
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout
    slide2.shapes.title.text = "Sales Performance Analysis"

    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Add formatted paragraphs
    p1 = tf.paragraphs[0]
    p1.text = "Executive Summary"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Our Q4 performance exceeded expectations with significant growth across all key metrics:"
    p2.font.size = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "• Revenue increased by 25% year-over-year"
    p3.level = 1

    p4 = tf.add_paragraph()
    p4.text = "• Customer acquisition grew by 40%"
    p4.level = 1

    p5 = tf.add_paragraph()
    p5.text = "• Market share expanded from 15% to 22%"
    p5.level = 1

    p6 = tf.add_paragraph()
    p6.text = "Key Success Factors"
    p6.font.bold = True
    p6.font.size = Pt(16)

    p7 = tf.add_paragraph()
    p7.text = "The remarkable growth can be attributed to our enhanced product offerings, strategic partnerships, and improved customer experience initiatives."

    # Add speaker notes
    slide2.notes_slide.notes_text_frame.text = (
        "This slide presents our Q4 sales performance overview. "
        "Key talking points: Emphasize the 25% revenue growth and 40% customer acquisition increase. "
        "Mention that market share expansion from 15% to 22% demonstrates strong competitive positioning. "
        "Be prepared to discuss the strategic initiatives that drove these results."
    )

    # Slide 3: Table with financial data
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide3.shapes.title.text = "Quarterly Financial Results"

    # Add table
    rows, cols = 5, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)

    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    # Set table headers
    headers = ["Quarter", "Revenue ($M)", "Profit ($M)", "Growth (%)"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    # Add data
    data = [
        ["Q1 2023", "45.2", "8.1", "12%"],
        ["Q2 2023", "52.8", "10.3", "18%"],
        ["Q3 2023", "58.9", "12.7", "15%"],
        ["Q4 2023", "67.3", "15.4", "25%"],
    ]

    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx, col_idx).text = cell_data

    # Add speaker notes
    slide3.notes_slide.notes_text_frame.text = (
        "This table shows our quarterly progression throughout 2023. "
        "Notice the consistent growth trend with Q4 showing the strongest performance. "
        "Revenue grew from $45.2M in Q1 to $67.3M in Q4, representing a 49% increase. "
        "Profit margins improved significantly, reaching $15.4M in Q4. "
        "The growth percentages show accelerating momentum, particularly in Q4 with 25% growth."
    )

    # Slide 4: Chart with sales data
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide4.shapes.title.text = "Monthly Sales Trends"

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    chart_data.add_series("Product A", (100, 125, 150, 175, 200, 225))
    chart_data.add_series("Product B", (80, 90, 110, 140, 160, 180))
    chart_data.add_series("Product C", (60, 75, 85, 95, 120, 140))

    # Add chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
    chart = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "Sales Performance by Product Line"

    # Add speaker notes
    slide4.notes_slide.notes_text_frame.text = (
        "This chart illustrates the monthly sales trends for our three main product lines. "
        "Product A shows the strongest performance with consistent growth from 100 to 225 units. "
        "Product B demonstrates steady improvement, reaching 180 units by June. "
        "Product C shows accelerating growth, particularly in the last two months. "
        "The overall trend indicates a healthy product portfolio with all lines contributing to growth."
    )

    # Slide 5: Mixed content slide
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout
    slide5.shapes.title.text = "Regional Performance Summary"

    # Add text content
    content = slide5.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "North America: Leading Market"
    p1.font.bold = True
    p1.font.size = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Generated $28.5M in revenue (42% of total), driven by strong enterprise adoption and new partnership agreements."

    p3 = tf.add_paragraph()
    p3.text = "Europe: Emerging Opportunities"
    p3.font.bold = True
    p3.font.size = Pt(16)

    p4 = tf.add_paragraph()
    p4.text = "Achieved $18.7M in revenue (28% of total) with significant growth in Germany and UK markets."

    # Add small table for regional breakdown
    left = Inches(1)
    top = Inches(4.5)
    width = Inches(6)
    height = Inches(1.5)

    small_table = slide5.shapes.add_table(4, 3, left, top, width, height).table

    # Headers
    headers = ["Region", "Revenue ($M)", "Market Share"]
    for i, header in enumerate(headers):
        small_table.cell(0, i).text = header
        small_table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    # Data
    regional_data = [
        ["North America", "28.5", "42%"],
        ["Europe", "18.7", "28%"],
        ["Asia Pacific", "20.1", "30%"],
    ]

    for row_idx, row_data in enumerate(regional_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            small_table.cell(row_idx, col_idx).text = cell_data

    # Add speaker notes
    slide5.notes_slide.notes_text_frame.text = (
        "This slide combines textual analysis with supporting data table. "
        "North America remains our strongest market, but note the balanced distribution across regions. "
        "Europe shows promising growth potential, especially in enterprise segments. "
        "Asia Pacific, while showing strong numbers, presents opportunities for expansion. "
        "The regional diversification reduces market concentration risk and provides multiple growth avenues."
    )

    # Slide 6: Future projections with line chart
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide6.shapes.title.text = "2024 Growth Projections"

    # Create line chart data
    line_chart_data = CategoryChartData()
    line_chart_data.categories = ["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024"]
    line_chart_data.add_series("Conservative", (70, 78, 85, 92))
    line_chart_data.add_series("Optimistic", (75, 85, 95, 108))
    line_chart_data.add_series("Stretch Goal", (80, 92, 105, 120))

    # Add line chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    line_chart = slide6.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, line_chart_data
    ).chart

    line_chart.has_title = True
    line_chart.chart_title.text_frame.text = "Revenue Projections ($M)"

    # Add speaker notes
    slide6.notes_slide.notes_text_frame.text = (
        "Our 2024 projections show three scenarios based on market conditions and execution capabilities. "
        "Conservative scenario assumes 15-20% growth, reaching $92M by Q4. "
        "Optimistic scenario projects 25-30% growth, achieving $108M in Q4. "
        "Stretch goal represents aggressive expansion with potential $120M Q4 revenue. "
        "We're targeting the optimistic scenario while preparing contingencies for the conservative case."
    )

    # Slide 7: Customer Demographics Table
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide7.shapes.title.text = "Customer Demographics Analysis"

    # Add demographics table
    rows, cols = 6, 5
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(3.5)

    demo_table = slide7.shapes.add_table(rows, cols, left, top, width, height).table

    # Set headers
    demo_headers = [
        "Age Group",
        "Percentage",
        "Revenue Share",
        "Growth Rate",
        "Retention",
    ]
    for i, header in enumerate(demo_headers):
        demo_table.cell(0, i).text = header
        demo_table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    # Add demographic data
    demo_data = [
        ["18-25", "15%", "12%", "45%", "78%"],
        ["26-35", "35%", "38%", "28%", "85%"],
        ["36-45", "28%", "32%", "15%", "92%"],
        ["46-55", "18%", "16%", "8%", "95%"],
        ["55+", "4%", "2%", "5%", "88%"],
    ]

    for row_idx, row_data in enumerate(demo_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            demo_table.cell(row_idx, col_idx).text = cell_data

    slide7.notes_slide.notes_text_frame.text = (
        "Customer demographics reveal interesting patterns in our user base. "
        "The 26-35 age group represents our largest segment at 35% of customers and 38% of revenue. "
        "Younger demographics (18-25) show highest growth at 45% but lower retention at 78%. "
        "Older segments demonstrate higher retention rates, with 46-55 age group at 95% retention. "
        "This data suggests opportunities for retention improvement in younger segments."
    )

    # Slide 8: Product Portfolio Pie Chart
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide8.shapes.title.text = "Product Portfolio Distribution"

    # Create pie chart data
    pie_chart_data = CategoryChartData()
    pie_chart_data.categories = [
        "Enterprise Software",
        "Mobile Apps",
        "Cloud Services",
        "Consulting",
        "Hardware",
    ]
    pie_chart_data.add_series("Revenue Share", (45, 25, 18, 8, 4))

    # Add pie chart
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    pie_chart = slide8.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, pie_chart_data
    ).chart

    pie_chart.has_title = True
    pie_chart.chart_title.text_frame.text = "Revenue by Product Category"

    slide8.notes_slide.notes_text_frame.text = (
        "Our product portfolio shows strong diversification across five key categories. "
        "Enterprise Software dominates with 45% of revenue, reflecting our B2B focus. "
        "Mobile Apps contribute 25%, showing strong consumer market presence. "
        "Cloud Services at 18% represent our fastest-growing segment. "
        "Consulting services provide 8% steady revenue with high margins. "
        "Hardware, while only 4%, offers strategic partnerships and ecosystem benefits."
    )

    # Slide 9: Competitive Analysis Matrix
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout
    slide9.shapes.title.text = "Competitive Landscape Assessment"

    # Add competitive analysis text
    content = slide9.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "Market Position Analysis"
    p1.font.bold = True
    p1.font.size = Pt(18)

    p2 = tf.add_paragraph()
    p2.text = "Competitive Advantages:"
    p2.font.bold = True
    p2.font.size = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "• Superior customer support with 24/7 availability"
    p3.level = 1

    p4 = tf.add_paragraph()
    p4.text = "• Advanced AI-driven analytics capabilities"
    p4.level = 1

    p5 = tf.add_paragraph()
    p5.text = "• Comprehensive integration ecosystem"
    p5.level = 1

    p6 = tf.add_paragraph()
    p6.text = "Areas for Improvement:"
    p6.font.bold = True
    p6.font.size = Pt(14)

    p7 = tf.add_paragraph()
    p7.text = "• Mobile platform feature parity"
    p7.level = 1

    p8 = tf.add_paragraph()
    p8.text = "• International market expansion"
    p8.level = 1

    # Add competitive matrix table
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(2)

    comp_table = slide9.shapes.add_table(4, 4, left, top, width, height).table

    # Headers
    comp_headers = ["Competitor", "Market Share", "Strengths", "Weaknesses"]
    for i, header in enumerate(comp_headers):
        comp_table.cell(0, i).text = header
        comp_table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    # Competitive data
    comp_data = [
        ["Company A", "35%", "Brand Recognition", "Limited Innovation"],
        ["Company B", "22%", "Cost Leadership", "Poor Support"],
        ["Our Company", "18%", "Technology Edge", "Market Penetration"],
    ]

    for row_idx, row_data in enumerate(comp_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            comp_table.cell(row_idx, col_idx).text = cell_data

    slide9.notes_slide.notes_text_frame.text = (
        "Competitive analysis reveals our strong technology position despite smaller market share. "
        "Company A leads with 35% share but lacks innovation velocity. "
        "Company B competes on price but suffers from support issues. "
        "Our 18% share is offset by superior technology and customer satisfaction. "
        "Focus should be on leveraging our tech advantages to gain market share."
    )

    # Slide 10: Financial KPIs Dashboard
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide10.shapes.title.text = "Key Performance Indicators Dashboard"

    # Add KPI table
    rows, cols = 8, 4
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4.5)

    kpi_table = slide10.shapes.add_table(rows, cols, left, top, width, height).table

    # KPI headers
    kpi_headers = ["Metric", "Current", "Target", "Status"]
    for i, header in enumerate(kpi_headers):
        kpi_table.cell(0, i).text = header
        kpi_table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    # KPI data
    kpi_data = [
        ["Monthly Recurring Revenue", "$5.2M", "$6.0M", "On Track"],
        ["Customer Acquisition Cost", "$450", "$400", "Needs Work"],
        ["Lifetime Value", "$2,800", "$3,000", "Good"],
        ["Churn Rate", "3.2%", "2.5%", "Improving"],
        ["Net Promoter Score", "68", "70", "Close"],
        ["Gross Margin", "72%", "75%", "Improving"],
        ["Employee Satisfaction", "8.1/10", "8.5/10", "Good"],
    ]

    for row_idx, row_data in enumerate(kpi_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            kpi_table.cell(row_idx, col_idx).text = cell_data

    slide10.notes_slide.notes_text_frame.text = (
        "Our KPI dashboard shows mixed but generally positive performance. "
        "MRR is tracking well toward $6M target, currently at $5.2M. "
        "CAC needs attention at $450, above our $400 target. "
        "LTV of $2,800 provides healthy unit economics with 6:1 LTV:CAC ratio. "
        "Churn improvement from 4.1% to 3.2% shows retention initiatives working. "
        "Focus areas: reduce CAC through channel optimization and improve NPS."
    )

    # Slide 11: Technology Roadmap
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout
    slide11.shapes.title.text = "Technology Roadmap 2024-2025"

    # Add roadmap content
    content = slide11.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "Q1 2024 Priorities"
    p1.font.bold = True
    p1.font.size = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "• AI-powered recommendation engine deployment"
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "• Mobile app performance optimization"
    p3.level = 1

    p4 = tf.add_paragraph()
    p4.text = "• Enhanced security framework implementation"
    p4.level = 1

    p5 = tf.add_paragraph()
    p5.text = "Q2-Q3 2024 Initiatives"
    p5.font.bold = True
    p5.font.size = Pt(16)

    p6 = tf.add_paragraph()
    p6.text = "• Real-time analytics platform launch"
    p6.level = 1

    p7 = tf.add_paragraph()
    p7.text = "• API ecosystem expansion"
    p7.level = 1

    p8 = tf.add_paragraph()
    p8.text = "• Multi-tenant architecture migration"
    p8.level = 1

    p9 = tf.add_paragraph()
    p9.text = "Q4 2024 & Beyond"
    p9.font.bold = True
    p9.font.size = Pt(16)

    p10 = tf.add_paragraph()
    p10.text = "• Machine learning automation suite"
    p10.level = 1

    p11 = tf.add_paragraph()
    p11.text = "• Global infrastructure expansion"
    p11.level = 1

    p12 = tf.add_paragraph()
    p12.text = "• Next-generation user interface rollout"
    p12.level = 1

    slide11.notes_slide.notes_text_frame.text = (
        "Our technology roadmap focuses on three key themes: intelligence, performance, and scale. "
        "Q1 priorities center on AI capabilities and mobile optimization for immediate user impact. "
        "Mid-year initiatives build platform capabilities for long-term competitive advantage. "
        "Q4 and beyond targets transformational capabilities including ML automation. "
        "Each initiative aligns with customer feedback and market opportunities."
    )

    # Slide 12: Summary and Next Steps
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout
    slide12.shapes.title.text = "Executive Summary & Action Items"

    # Add summary content
    content = slide12.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "Key Achievements"
    p1.font.bold = True
    p1.font.size = Pt(18)

    p2 = tf.add_paragraph()
    p2.text = "✓ 25% revenue growth exceeding targets"
    p2.font.size = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "✓ Successful market share expansion to 22%"
    p3.font.size = Pt(14)

    p4 = tf.add_paragraph()
    p4.text = "✓ Strong customer retention improvements"
    p4.font.size = Pt(14)

    p5 = tf.add_paragraph()
    p5.text = "Immediate Action Items"
    p5.font.bold = True
    p5.font.size = Pt(18)

    p6 = tf.add_paragraph()
    p6.text = "1. Optimize customer acquisition costs"
    p6.level = 1
    p6.font.size = Pt(14)

    p7 = tf.add_paragraph()
    p7.text = "2. Accelerate AI feature development"
    p7.level = 1
    p7.font.size = Pt(14)

    p8 = tf.add_paragraph()
    p8.text = "3. Expand European market presence"
    p8.level = 1
    p8.font.size = Pt(14)

    p9 = tf.add_paragraph()
    p9.text = "4. Enhance mobile platform capabilities"
    p9.level = 1
    p9.font.size = Pt(14)

    # Add final summary table
    left = Inches(1)
    top = Inches(4.5)
    width = Inches(8)
    height = Inches(1.5)

    summary_table = slide12.shapes.add_table(4, 3, left, top, width, height).table

    # Summary headers
    summary_headers = ["Priority", "Owner", "Timeline"]
    for i, header in enumerate(summary_headers):
        summary_table.cell(0, i).text = header
        summary_table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    # Summary data
    summary_data = [
        ["CAC Optimization", "Marketing Team", "Q1 2024"],
        ["AI Development", "Engineering Team", "Q2 2024"],
        ["European Expansion", "Sales Team", "Q1-Q2 2024"],
    ]

    for row_idx, row_data in enumerate(summary_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            summary_table.cell(row_idx, col_idx).text = cell_data

    slide12.notes_slide.notes_text_frame.text = (
        "This final slide summarizes our key achievements and establishes clear action items. "
        "The 25% revenue growth demonstrates strong execution of our strategy. "
        "Market share expansion to 22% positions us well for continued growth. "
        "Action items are prioritized based on impact and feasibility. "
        "CAC optimization is critical for sustainable growth economics. "
        "AI development maintains our competitive technology advantage. "
        "European expansion diversifies revenue streams and reduces market risk."
    )

    # Save presentation
    prs.save(filename)
    print(f"✅ Created comprehensive test presentation: {filename}")
    print("\n📋 12-Slide Presentation Contents:")
    print("   • Slide 1: Title slide with speaker notes")
    print("   • Slide 2: Rich formatted text with bullet points")
    print("   • Slide 3: Financial data table (4x4)")
    print("   • Slide 4: Column chart with 3 data series")
    print("   • Slide 5: Mixed content (text + table)")
    print("   • Slide 6: Line chart with projections")
    print("   • Slide 7: Customer demographics table (5x5)")
    print("   • Slide 8: Product portfolio pie chart")
    print("   • Slide 9: Competitive analysis matrix (mixed content)")
    print("   • Slide 10: Financial KPIs dashboard (7x4 table)")
    print("   • Slide 11: Technology roadmap (structured text)")
    print("   • Slide 12: Summary and action items table")
    print("\n🎯 Perfect for testing:")
    print("   ✓ Text extraction and formatting preservation")
    print("   ✓ Table data extraction with complete content")
    print("   ✓ Chart metadata and series data extraction")
    print("   ✓ Speaker notes extraction for all slides")
    print("   ✓ Mixed content handling (text + tables + charts)")
    print("   ✓ Multithreading with 12 slides for concurrency testing")
    print("   ✓ Various chart types (column, line, pie)")
    print("   ✓ Different table sizes and structures")
    print("   ✓ Rich text formatting with hierarchical content")
    print("   ✓ Comprehensive speaker notes for LLM consolidation")
    print("\n🚀 Concurrency Testing Features:")
    print(f"   • 12 slides perfect for batch_size testing")
    print(f"   • Multiple content types per slide")
    print(f"   • Rich metadata for comprehensive extraction")
    print(f"   • Suitable for rate limit testing with LLM consolidation")

    return filename


if __name__ == "__main__":
    create_comprehensive_test_presentation()
